#!/usr/bin/env python3
"""
Server-side PPTX and PDF generator for notso.ai proposal system.
Reads proposal JSON from stdin, outputs file to specified path.

Usage:
  echo '{"format":"pptx","proposal":{...},"client":{...}}' | python3 generate.py
"""
import sys, json, os, re
from pathlib import Path

# ─── PPTX Generation ───
def generate_pptx(proposal, client, output_path, selected_slides=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    W, H = Inches(10), Inches(5.625)
    M = Inches(0.6)

    # ─── Brand palette (4 user-chosen colors) ───
    # color1 = main brand (primary headers, hero backgrounds)
    # color2 = accent 1 (borders, secondary highlights)
    # color3 = accent 2 (pills, tertiary highlights — often yellow/orange)
    # color4 = accent 3 (tags, callouts — often red/coral)
    c1_hex = (client.get('color1') or '#0B3C8C').lstrip('#')
    c2_hex = (client.get('color2') or '#E63946').lstrip('#')
    c3_hex = (client.get('color3') or '#F5A623').lstrip('#')
    c4_hex = (client.get('color4') or '#E74C3C').lstrip('#')

    def hex_to_rgb(h):
        h = h.lstrip('#')
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    def luminance(h):
        h = h.lstrip('#')
        r,g,b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
        return 0.299*r + 0.587*g + 0.114*b

    def darken(h, pct):
        r = max(0, int(int(h[0:2],16) * (1-pct)))
        g = max(0, int(int(h[2:4],16) * (1-pct)))
        b = max(0, int(int(h[4:6],16) * (1-pct)))
        return f'{r:02x}{g:02x}{b:02x}'

    def tint(h, pct):
        """Blend a hex color with white by pct (0..1)."""
        h = h.lstrip('#')
        r = int(h[0:2],16); g = int(h[2:4],16); b = int(h[4:6],16)
        r = int(r + (255 - r) * pct)
        g = int(g + (255 - g) * pct)
        b = int(b + (255 - b) * pct)
        return RGBColor(r, g, b)

    C1 = hex_to_rgb(c1_hex)                 # primary (main)
    C2 = hex_to_rgb(c2_hex)                 # accent 1
    C3 = hex_to_rgb(c3_hex)                 # accent 2
    C4 = hex_to_rgb(c4_hex)                 # accent 3
    ACCENT = C2                             # back-compat alias (ACCENT == accent 1)
    on_c1 = RGBColor(10,10,10) if luminance(c1_hex) > 160 else RGBColor(255,255,255)
    BLACK = RGBColor(26,26,26)              # slightly softer than pure black
    GRAY = RGBColor(107,114,128)
    GRAY_LIGHT = RGBColor(156,163,175)
    BORDER = RGBColor(229,231,235)          # card border
    WHITE = RGBColor(255,255,255)
    LIGHT_BG = RGBColor(249,250,251)        # cooler near-white
    CARD_BG = RGBColor(255,255,255)
    # Tinted pastel variants for emotion cards + soft backgrounds
    C1_TINT = tint(c1_hex, 0.78)            # light blue (mascot slot bg)
    C2_TINT = tint(c2_hex, 0.82)
    C3_TINT = tint(c3_hex, 0.78)
    C4_TINT = tint(c4_hex, 0.82)
    MASCOT_SLOT = C1_TINT                   # light blue rectangle (mascot placement convention)
    # Hills + clouds illustration colors
    SKY_BLUE = RGBColor(179,229,252)        # sky
    HILL_GREEN_1 = RGBColor(139,195,74)     # foreground hill
    HILL_GREEN_2 = RGBColor(104,159,56)     # back hill (darker)
    CLOUD_WHITE = RGBColor(255,255,255)

    # ─── Typography: Poppins (matches notso.ai reference decks) ───
    # Falls back to Calibri / Arial gracefully on systems without Poppins installed
    FONT_HEAD = 'Poppins'
    FONT_BODY = 'Poppins'

    def strip_emoji(s):
        if not s: return ''
        return re.sub(r'[\U0001F300-\U0001FAFF\U00002600-\U000027BF\U0000FE00-\U0000FE0F\U0001F900-\U0001F9FF\U0000200D\U000020E3\U000E0020-\U000E007F]', '', str(s)).strip()

    def add_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_rect(slide, left, top, width, height, fill_color):
        from pptx.util import Emu as E
        shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def add_text_box(slide, left, top, width, height, text, font_size=11, bold=False, italic=False, color=BLACK, align=PP_ALIGN.LEFT, font_name=None):
        if font_name is None:
            font_name = FONT_BODY
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = strip_emoji(text)
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
        return txBox

    def add_footer(slide):
        # Brand footer: notso.ai logo (accent color) in bottom-right
        add_text_box(slide, Inches(8.5), Inches(5.28), Inches(1.2), Inches(0.3),
                     'notso.ai', font_size=10, bold=True, color=ACCENT,
                     align=PP_ALIGN.RIGHT, font_name=FONT_HEAD)

    import datetime
    date_str = datetime.datetime.now().strftime('%d %B %Y')

    # Selected slides filter
    sel = set(selected_slides) if selected_slides else {'s1','s2','s3','s4','s5','s6','s7','s8','s9','s10','s11','s12','s13','s14','s15','s16','s17','s18'}

    # ── Helper: add image frame (placeholder or real image) ──
    PLACEHOLDER_BG = RGBColor(230, 233, 236)  # light neutral gray
    DARK_CARD = RGBColor(51, 51, 51)  # dark card background for pricing

    # Mascot images dict: { 'option_a': '/path/to/img.png', 'option_b': '...', 'cover': '...' }
    mascot_images = proposal.get('_mascot_images', {})

    def add_placeholder_img(slide, left, top, width, height, image_key=None, label='Insert Image'):
        """
        Image placement zone. If mascot_images[image_key] is a valid file path, insert
        the image (scaled to fit). Otherwise render a light-gray rounded-rectangle
        placeholder with a dashed border and centered label.
        """
        img_path = mascot_images.get(image_key, '') if image_key else ''

        if img_path and os.path.exists(img_path):
            try:
                return slide.shapes.add_picture(img_path, left, top, width, height)
            except Exception:
                pass  # fall through to placeholder

        # Rounded placeholder box with dashed border
        shape = add_rounded_rect(slide, left, top, width, height, PLACEHOLDER_BG)
        shape.line.color.rgb = RGBColor(180, 184, 190)
        shape.line.width = Pt(1.25)
        try:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass

        # Normalize dims to EMU for safe arithmetic (Length is int subclass)
        import pptx.util as _u
        L = int(left); T = int(top); W = int(width); H = int(height)

        # Only draw label if the zone is tall enough (>= 0.35")
        if H >= int(_u.Inches(0.35)):
            # Scale font with shorter dim; clamp 8..15
            shorter_in = min(W, H) / _u.Inches(1)
            fs = max(8, min(15, int(shorter_in * 7.5)))
            text_h = _u.Inches(0.42)
            text_top = T + (H - text_h) // 2
            add_text_box(slide, Emu(L), Emu(text_top), Emu(W), Emu(text_h),
                         label, font_size=fs, color=RGBColor(140, 145, 155),
                         align=PP_ALIGN.CENTER)
        return shape

    def add_rounded_rect(slide, left, top, width, height, fill_color, corner_radius=Inches(0.08)):
        """Add a rounded rectangle using MSO_SHAPE.ROUNDED_RECTANGLE (id=5)."""
        shape = slide.shapes.add_shape(5, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def add_oval(slide, left, top, width, height, fill_color):
        """Add an oval/circle shape (MSO_SHAPE.OVAL = 9)."""
        shape = slide.shapes.add_shape(9, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def add_pill(slide, left, top, width, height, text, fill_color, text_color, font_size=8, bold=True):
        """Small rounded pill badge (for tags, categories)."""
        shape = add_rounded_rect(slide, left, top, width, height, fill_color)
        add_text_box(slide, left, top + Emu(int(height * 0.15)), width, height - Emu(int(height * 0.2)),
                     text, font_size=font_size, bold=bold, color=text_color, align=PP_ALIGN.CENTER)
        return shape

    def add_title_lead(slide, headline, lead='', top=0.3, left=0.55, width=9.0,
                       title_size=26, title_color=None, title_bold=True,
                       lead_color=None, lead_size=11):
        """
        Render a page headline with a 1-sentence lead-text subtitle beneath it.
        Used on nearly every slide to ensure titles follow the 'sign on the door →
        welcoming words inside' pattern the client asked for.

        Returns the y-position (in inches) immediately below the lead so the caller
        can position its content directly after.
        """
        if title_color is None:
            title_color = C1
        if lead_color is None:
            lead_color = GRAY
        add_text_box(slide, Inches(left), Inches(top), Inches(width), Inches(0.65),
                     strip_emoji(headline or ''),
                     font_size=title_size, bold=title_bold, color=title_color, font_name=FONT_HEAD)
        lead_text = strip_emoji(lead or '')
        if lead_text:
            add_text_box(slide, Inches(left), Inches(top + 0.68), Inches(width), Inches(0.42),
                         lead_text, font_size=lead_size, italic=True, color=lead_color, font_name=FONT_BODY)
            return top + 1.18
        return top + 0.78

    def add_hero_band(slide, headline, lead='', band_color=None, title_size=24,
                      band_height=1.25, text_pad_x=0.55):
        """
        Notso reference style: Large black title at top-left with optional lead text.
        A thin brand-color accent line separates the title from content below.

        Returns the y-position immediately below the band so the caller can
        continue laying out content.
        """
        if band_color is None:
            band_color = C1
        # Large title (black, bold, Poppins — matching Fin 45pt style)
        add_text_box(slide, Inches(text_pad_x), Inches(0.3), Inches(10 - 2 * text_pad_x), Inches(0.7),
                     strip_emoji(headline or ''),
                     font_size=title_size, bold=True, color=BLACK, font_name=FONT_HEAD)
        # Lead subtitle in gray
        lead_text = strip_emoji(lead or '')
        if lead_text:
            add_text_box(slide, Inches(text_pad_x), Inches(0.95), Inches(10 - 2 * text_pad_x), Inches(0.35),
                         lead_text, font_size=11, italic=False, color=GRAY,
                         font_name=FONT_BODY)
        # Thin brand-color accent line
        line_y = 1.35 if lead_text else 1.05
        add_rect(slide, Inches(text_pad_x), Inches(line_y), Inches(2.0), Inches(0.04), band_color)
        return line_y + 0.2

    # ══════ SLIDE 1: COVER ══════
    # Notso reference style (Fin/Jumbo):
    #   - Light gray #F4F4F3 background
    #   - Left side: large bold title (black, 36pt), tagline below, date + "online" pill
    #   - Right side: large mascot image (hero placement)
    #   - Clean, minimal — no heavy colored bars
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    COVER_BG = RGBColor(244, 244, 243)  # warm off-white (#F4F4F3 from Fin ref)
    add_bg(s1, COVER_BG)
    d1 = proposal.get('s1', {})
    mascot_name = strip_emoji(d1.get('mascot_name') or proposal.get('mascot_name') or '')
    client_name = client.get('name', '')

    # ── "online" status pill (top-left) ──
    pill_y = Inches(0.4)
    add_rounded_rect(s1, Inches(0.6), pill_y, Inches(1.2), Inches(0.32), WHITE)
    add_text_box(s1, Inches(0.7), pill_y + Inches(0.04), Inches(1.0), Inches(0.25),
                 'online', font_size=9, color=BLACK, font_name=FONT_BODY)

    # ── Main title: "Mascot – De visuele ai-agent voor CLIENT" ──
    # Text stays in left 45% (0.6 to 4.8"), mascot in right 45% (5.5 to 10")
    _lead1 = strip_emoji(d1.get('lead', '') or d1.get('greeting', '') or d1.get('tagline', ''))
    title_line = f'{mascot_name} – ' + (_lead1 if _lead1 else f'The AI mascot for {client_name}')
    # Truncate title if too long (max ~80 chars to fit in box without overflow)
    if len(title_line) > 80:
        title_line = title_line[:77] + '...'
    add_text_box(s1, Inches(0.6), Inches(1.0), Inches(4.2), Inches(2.8),
                 title_line, font_size=32, bold=False, color=BLACK, font_name=FONT_HEAD)

    # ── Italic sub-pill: "Mascot // online" ──
    add_text_box(s1, Inches(0.6), Inches(4.0), Inches(3.0), Inches(0.3),
                 f'{mascot_name} // online', font_size=9, italic=True, color=BLACK, font_name=FONT_BODY)

    # ── Date validity line ──
    add_text_box(s1, Inches(0.6), Inches(4.4), Inches(3.5), Inches(0.3),
                 date_str, font_size=11, color=BLACK, font_name=FONT_BODY)

    # ── Mascot hero image (right ~45% of slide, no overlap with text) ──
    add_placeholder_img(s1, Inches(5.5), Inches(0.0), Inches(4.5), Inches(5.2),
                        image_key='cover', label='Mascot')

    # ── Small brand accent strip at very bottom ──
    add_rect(s1, Inches(0), Inches(5.35), W, Inches(0.275), C1)

    # ══════ SLIDE 3: PAIN POINTS ══════
    # Layout-03 reference: full-width C1 HERO BAND containing title + lead in white,
    # then 3 numbered cards with top C1 bar, large 01/02/03 number, bold title, gray desc.
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s2, LIGHT_BG)  # light gray background like reference
    d2 = proposal.get('s3', {})

    after_title = add_hero_band(s2,
        d2.get('headline', 'Pain Points'),
        d2.get('lead', '') or d2.get('intro', ''),
        band_color=C1, title_size=22, band_height=1.25)

    # 3 numbered cards
    points = d2.get('points', [])
    n_points = min(len(points), 3) or 1
    card_gap = 0.35
    card_w = (8.9 - card_gap * (n_points - 1)) / n_points
    card_h = 3.3
    card_top = max(1.75, after_title + 0.25)

    for i, p in enumerate(points[:3]):
        cx = 0.55 + i * (card_w + card_gap)
        # White card (no border) with a C1 top bar
        card_shape = add_rect(s2, Inches(cx), Inches(card_top), Inches(card_w), Inches(card_h), CARD_BG)
        card_shape.line.color.rgb = BORDER
        card_shape.line.width = Pt(0.5)
        # Top C1 bar
        add_rect(s2, Inches(cx), Inches(card_top), Inches(card_w), Inches(0.08), C1)

        # Number (large, C1)
        add_text_box(s2, Inches(cx + 0.3), Inches(card_top + 0.3), Inches(card_w - 0.6), Inches(0.7),
                     str(i+1).zfill(2), font_size=36, bold=True, color=C1, font_name=FONT_HEAD)
        # Title (bold black)
        add_text_box(s2, Inches(cx + 0.3), Inches(card_top + 1.15), Inches(card_w - 0.6), Inches(0.5),
                     strip_emoji(p.get('title','')), font_size=14, bold=True, color=BLACK, font_name=FONT_HEAD)
        # Description (gray)
        add_text_box(s2, Inches(cx + 0.3), Inches(card_top + 1.7), Inches(card_w - 0.6), Inches(card_h - 1.85),
                     strip_emoji(p.get('desc','')), font_size=10, color=GRAY)
    add_footer(s2)

    # ══════ SLIDE 5: CORE FEATURES (text-only 2x2) ══════
    # Layout-05 reference: light gray background on title area + lead, then 2x2
    # text-only cards. Each card has a small C1 underline beneath the title, a
    # bold title, and a gray description. NO image placeholders.
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s3, WHITE)
    d3 = proposal.get('s5', {})

    after_title = add_title_lead(s3,
        d3.get('headline', 'Core Features'),
        d3.get('lead', '') or d3.get('intro', ''),
        top=0.35, title_size=24)

    feats = d3.get('features', [])
    card_w = 4.35
    card_h = 1.65
    gap_x = 0.25
    gap_y = 0.25
    grid_top = max(1.65, after_title + 0.15)

    # Cycle through all 4 brand colors for the card accents
    accent_cycle = [C1, C3, C4, C2]
    # Each card: text on left, light-blue mascot slot on right (per layout-05 reference)
    mascot_slot_w = 1.35
    text_w = card_w - mascot_slot_w - 0.5
    for i, f in enumerate(feats[:4]):
        col, row = i % 2, i // 2
        fx = 0.55 + col * (card_w + gap_x)
        fy = grid_top + row * (card_h + gap_y)
        accent = accent_cycle[i % len(accent_cycle)]

        # White card with thin border
        card_shape = add_rect(s3, Inches(fx), Inches(fy), Inches(card_w), Inches(card_h), CARD_BG)
        card_shape.line.color.rgb = BORDER
        card_shape.line.width = Pt(0.75)
        # Colored left vertical accent bar — cycles through all 4 brand colors
        add_rect(s3, Inches(fx), Inches(fy), Inches(0.1), Inches(card_h), accent)

        # Bold title (C1)
        add_text_box(s3, Inches(fx + 0.3), Inches(fy + 0.2), Inches(text_w), Inches(0.4),
                     strip_emoji(f.get('title','')),
                     font_size=14, bold=True, color=C1, font_name=FONT_HEAD)
        # Colored underline bar in the card's accent color
        add_rect(s3, Inches(fx + 0.3), Inches(fy + 0.62), Inches(0.6), Inches(0.035), accent)
        # Description (gray) — constrained to left area
        add_text_box(s3, Inches(fx + 0.3), Inches(fy + 0.75), Inches(text_w), Inches(card_h - 0.9),
                     strip_emoji(f.get('desc','')), font_size=9, color=GRAY)

        # Light-blue mascot slot on the right of the card
        slot_x = fx + card_w - mascot_slot_w - 0.18
        slot_y = fy + 0.18
        slot_h = card_h - 0.36
        mslot = add_rect(s3, Inches(slot_x), Inches(slot_y),
                         Inches(mascot_slot_w), Inches(slot_h), C1_TINT)
        mslot.line.color.rgb = C1
        mslot.line.width = Pt(0.75)
        try:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            mslot.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass

    add_footer(s3)

    # ══════ SLIDE 6: MASCOT SELECTION ══════
    # Layout-06 reference: two side-by-side cards, each with "Option A/B" small
    # dark label at top, a C1 underline, a huge light-blue mascot slot, the name
    # below (centered, bold black), then description below.
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s4, WHITE)
    d4 = proposal.get('s6', {})

    after_title = add_title_lead(s4,
        d4.get('headline', 'Mascot Design'),
        d4.get('lead', ''),
        top=0.35, title_size=26)

    opts = [d4.get('option_a'), d4.get('option_b'), d4.get('option_c')]
    opts = [o for o in opts if o]
    if not opts:
        opts = d4.get('options', [])
    num_opts = max(min(len(opts), 3), 1)
    col_gap = 0.4
    col_w = (8.9 - col_gap * (num_opts - 1)) / num_opts
    card_top = max(1.4, after_title + 0.15)
    card_h = 5.625 - card_top - 0.35

    for i, opt in enumerate(opts[:num_opts]):
        ox = 0.55 + i * (col_w + col_gap)

        # White card with thin border
        card_shape = add_rect(s4, Inches(ox), Inches(card_top), Inches(col_w), Inches(card_h), CARD_BG)
        card_shape.line.color.rgb = BORDER
        card_shape.line.width = Pt(0.75)

        # Option A/B label (dark, small)
        add_text_box(s4, Inches(ox + 0.3), Inches(card_top + 0.18), Inches(col_w - 0.6), Inches(0.35),
                     'Option ' + chr(65 + i), font_size=11, bold=True, color=BLACK, font_name=FONT_HEAD)
        # C1 underline bar
        add_rect(s4, Inches(ox + 0.3), Inches(card_top + 0.55), Inches(col_w - 0.6), Inches(0.03), C1)

        # Large light-blue mascot slot
        slot_top = card_top + 0.75
        slot_h = 2.2
        option_key = ['option_a', 'option_b', 'option_c'][i] if i < 3 else None
        add_placeholder_img(s4, Inches(ox + 0.3), Inches(slot_top),
                            Inches(col_w - 0.6), Inches(slot_h),
                            image_key=option_key, label='Mascot')

        # Mascot name (centered, bold black)
        name_y = slot_top + slot_h + 0.15
        add_text_box(s4, Inches(ox), Inches(name_y), Inches(col_w), Inches(0.45),
                     strip_emoji(opt.get('name','')),
                     font_size=18, bold=True, color=BLACK,
                     align=PP_ALIGN.CENTER, font_name=FONT_HEAD)

        # Description (gray)
        desc_y = name_y + 0.55
        add_text_box(s4, Inches(ox + 0.3), Inches(desc_y), Inches(col_w - 0.6),
                     Inches(card_top + card_h - desc_y - 0.15),
                     strip_emoji(opt.get('desc','')), font_size=10, color=GRAY)

    add_footer(s4)

    # ══════ SLIDE 7: MASCOT DESIGN (LALA-style) ══════
    # Layout-07 reference:
    #   - HUGE mascot name upper-left
    #   - Short personality description beneath
    #   - "Personality" label followed by 3 tone bars (Formal/Informal etc.)
    #     where the filled portion is C3 (accent 2 — yellow) on a dark track
    #   - Right side: big light-blue mascot slot with 2 speech bubbles floating
    #   - Catchphrases bullet list bottom-left, tone description bottom-right
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s5, WHITE)
    d5 = proposal.get('s7', {})
    mascot = strip_emoji(d5.get('name', d1.get('mascot_name', '')))
    _lead5 = strip_emoji(d5.get('lead', ''))

    # ── LEFT PANEL ──
    # HUGE mascot name (C1)
    add_text_box(s5, Inches(0.55), Inches(0.3), Inches(4.0), Inches(0.95),
                 mascot.upper() if mascot else '',
                 font_size=40, bold=True, color=C1, font_name=FONT_HEAD)
    # Lead-text sitting directly below the name
    if _lead5:
        add_text_box(s5, Inches(0.55), Inches(1.25), Inches(3.9), Inches(0.5),
                     _lead5, font_size=11, italic=True, color=GRAY, font_name=FONT_BODY)
    # Personality description
    add_text_box(s5, Inches(0.55), Inches(1.75), Inches(3.9), Inches(1.0),
                 strip_emoji(d5.get('personality','')), font_size=10, color=BLACK)

    # Personality section label
    add_text_box(s5, Inches(0.55), Inches(2.85), Inches(2.0), Inches(0.35),
                 'Personality', font_size=13, bold=True, color=BLACK, font_name=FONT_HEAD)

    # 3 tone bars (dark track, C3 fill)
    bar_labels = [('Formal', 'Informal'), ('Concise', 'Detailed'), ('Serious', 'Playful')]
    for bi, (left_l, right_l) in enumerate(bar_labels):
        by = 3.25 + bi * 0.42
        add_text_box(s5, Inches(0.55), Inches(by), Inches(1.0), Inches(0.2),
                     left_l, font_size=8, color=GRAY)
        add_text_box(s5, Inches(2.25), Inches(by), Inches(1.0), Inches(0.2),
                     right_l, font_size=8, color=GRAY, align=PP_ALIGN.RIGHT)
        # Dark track
        add_rect(s5, Inches(0.55), Inches(by + 0.22), Inches(2.7), Inches(0.1), BLACK)
        # C3 fill (accent 2 — yellow/orange)
        fill_w = 1.2 + (bi * 0.25)
        add_rect(s5, Inches(0.55), Inches(by + 0.22), Inches(fill_w), Inches(0.1), C3)

    # ── RIGHT PANEL: light-blue mascot slot ──
    add_placeholder_img(s5, Inches(4.7), Inches(0.3), Inches(5.0), Inches(3.0),
                        image_key='cover', label='Mascot')

    # Speech bubbles floating on the right
    phrases = d5.get('phrases', [])
    if phrases and len(phrases) >= 1:
        add_rounded_rect(s5, Inches(7.45), Inches(0.5), Inches(2.1), Inches(0.75), WHITE)
        add_text_box(s5, Inches(7.55), Inches(0.62), Inches(1.95), Inches(0.55),
                     '"' + strip_emoji(phrases[0]) + '"', font_size=8, italic=True, color=BLACK)
    if phrases and len(phrases) >= 2:
        add_rounded_rect(s5, Inches(4.8), Inches(2.1), Inches(2.1), Inches(0.75), WHITE)
        add_text_box(s5, Inches(4.9), Inches(2.22), Inches(1.95), Inches(0.55),
                     '"' + strip_emoji(phrases[1]) + '"', font_size=8, italic=True, color=BLACK)
    if phrases and len(phrases) >= 3:
        add_rounded_rect(s5, Inches(7.45), Inches(2.7), Inches(2.1), Inches(0.75), WHITE)
        add_text_box(s5, Inches(7.55), Inches(2.82), Inches(1.95), Inches(0.55),
                     '"' + strip_emoji(phrases[2]) + '"', font_size=8, italic=True, color=BLACK)

    # Catchphrases (bottom-left)
    add_text_box(s5, Inches(0.55), Inches(4.55), Inches(3.5), Inches(0.3),
                 'Catchphrases', font_size=11, bold=True, color=BLACK, font_name=FONT_HEAD)
    if phrases:
        phrases_text = '\n'.join(['• ' + strip_emoji(p) for p in phrases[:3]])
        add_text_box(s5, Inches(0.55), Inches(4.82), Inches(4.0), Inches(0.85),
                     phrases_text, font_size=8, color=GRAY)

    # Tone description (bottom-right, gray italic on white)
    tone_desc = strip_emoji(d5.get('tone_desc', ''))
    if tone_desc:
        add_text_box(s5, Inches(4.8), Inches(4.55), Inches(4.9), Inches(0.85),
                     'Tone: ' + tone_desc, font_size=9, italic=True, color=GRAY)

    add_footer(s5)

    # ══════ SLIDE 9: CHAT DEMO ══════
    # Layout-09 reference: "Chat Demo: Name" title, left side has light-blue
    # mascot slot with tone text below, right side has chat window with C1 header
    # bar and alternating bubbles (C1 for bot, gray for user).
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s6, WHITE)
    d6 = proposal.get('s9', {})
    bot_name = strip_emoji(d6.get('mascot_name', mascot or 'Bot'))

    # Constrain to left side so lead subtitle doesn't hide behind chat window on right
    after_title = add_title_lead(s6,
        'Chat Demo: ' + bot_name,
        d6.get('lead', ''),
        top=0.35, width=3.3, title_size=22)

    # Left: light-blue mascot slot
    add_placeholder_img(s6, Inches(0.55), Inches(max(1.5, after_title + 0.1)),
                        Inches(3.2), Inches(3.2),
                        image_key='cover', label='Mascot')

    # Tone description below mascot (only if present)
    _td = strip_emoji(d5.get('tone_desc', '') if isinstance(d5, dict) else '')
    if _td:
        add_text_box(s6, Inches(0.55), Inches(4.45), Inches(3.2), Inches(1.0),
                     _td, font_size=9, italic=True, color=GRAY)

    # ── Right: Chat window mockup (polished card) ──
    chat_left = 3.95
    chat_top = 1.1
    chat_w = 5.75
    chat_h = 4.2

    # Chat window with border
    chat_bg = add_rect(s6, Inches(chat_left), Inches(chat_top), Inches(chat_w), Inches(chat_h), WHITE)
    chat_bg.line.color.rgb = BORDER
    chat_bg.line.width = Pt(0.75)

    # Chat header bar with mascot name
    add_rect(s6, Inches(chat_left), Inches(chat_top), Inches(chat_w), Inches(0.5), LIGHT_BG)
    add_text_box(s6, Inches(chat_left + 0.2), Inches(chat_top + 0.1), Inches(2), Inches(0.35),
                 bot_name, font_size=14, bold=True, color=C1, font_name=FONT_HEAD)
    # Close button (x)
    add_text_box(s6, Inches(chat_left + chat_w - 0.5), Inches(chat_top + 0.1), Inches(0.35), Inches(0.3),
                 'x', font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

    # Chat messages
    chat = d6.get('chat', [])
    cy = chat_top + 0.65
    max_msg_w = chat_w - 0.8
    bubble_padding = 0.12

    for msg in chat[:4]:
        role = msg.get('r') or msg.get('who') or msg.get('role') or ''
        is_user = role == 'user'
        msg_text = strip_emoji(msg.get('m') or msg.get('text') or msg.get('message') or '')
        label = 'User:  ' if is_user else bot_name + ':  '

        # Calculate bubble height based on text length
        chars_per_line = 38
        lines = max(1, len(msg_text) / chars_per_line + 0.5)
        mh = max(0.5, min(1.2, lines * 0.22 + 0.2))
        bw = min(max_msg_w, max(2.5, len(msg_text) / chars_per_line * max_msg_w * 0.5 + 1.5))

        if is_user:
            # User: gray bubble, right-aligned
            bx = chat_left + chat_w - 0.3 - bw
            bc = LIGHT_BG
            tc = BLACK
        else:
            # Bot: primary color bubble, left-aligned
            bx = chat_left + 0.3
            bc = C1
            tc = WHITE

        add_rounded_rect(s6, Inches(bx), Inches(cy), Inches(bw), Inches(mh), bc)
        add_text_box(s6, Inches(bx + bubble_padding), Inches(cy + 0.08),
                     Inches(bw - bubble_padding * 2), Inches(mh - 0.16),
                     label + msg_text, font_size=9, color=tc)
        cy += mh + 0.1

        if cy > chat_top + chat_h - 0.7:
            break

    # Input bar at bottom (accent color)
    input_y = chat_top + chat_h - 0.5
    add_rounded_rect(s6, Inches(chat_left + 0.3), Inches(input_y), Inches(chat_w - 0.6), Inches(0.35), ACCENT)
    add_text_box(s6, Inches(chat_left + 0.45), Inches(input_y + 0.07), Inches(3), Inches(0.25),
                 'Say something...', font_size=9, italic=True, color=WHITE)

    add_footer(s6)

    # ══════ SLIDE 12: DATA & INSIGHTS ══════
    # Layout-12 reference:
    #   - Headline = pipe-separated badges (e.g. "即時數據儀表板 | 情感分析 | ...")
    #   - Lead-text under the title
    #   - Left: DASHBOARD SCREENSHOT PLACEHOLDER (dashed border, 16:10, gray)
    #     → user will paste their own screenshot here
    #   - Right: "Standard Included" label, intro text, bulleted metric list
    #     with a C4 bullet dot, then the same badges repeated at the bottom
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s7, WHITE)
    d7 = proposal.get('s12', {})

    # Title: use pipe-separated badges if present, else headline
    badges = d7.get('badges', []) or []
    headline_text = (' | '.join([strip_emoji(b) for b in badges])
                     if badges else strip_emoji(d7.get('headline', 'Data & Insights')))
    add_text_box(s7, Inches(0.55), Inches(0.35), Inches(9.0), Inches(0.55),
                 headline_text, font_size=18, bold=True, color=C1, font_name=FONT_HEAD)
    # Lead-text
    _lead7 = strip_emoji(d7.get('lead', ''))
    if _lead7:
        add_text_box(s7, Inches(0.55), Inches(0.88), Inches(9.0), Inches(0.4),
                     _lead7, font_size=11, italic=True, color=GRAY, font_name=FONT_BODY)

    # ── LEFT: Dashboard screenshot placeholder (dashed gray box, 16:10) ──
    dash_x, dash_y = Inches(0.55), Inches(1.5)
    dash_w, dash_h = Inches(5.2), Inches(3.25)
    dash_card = add_rect(s7, dash_x, dash_y, dash_w, dash_h, RGBColor(243, 244, 246))
    dash_card.line.color.rgb = RGBColor(156, 163, 175)
    dash_card.line.width = Pt(1.25)
    try:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        dash_card.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    except Exception:
        pass
    add_text_box(s7, dash_x, dash_y + Inches(1.35), dash_w, Inches(0.35),
                 'Dashboard Screenshot',
                 font_size=12, bold=True, color=RGBColor(107, 114, 128),
                 align=PP_ALIGN.CENTER, font_name=FONT_HEAD)
    add_text_box(s7, dash_x, dash_y + Inches(1.7), dash_w, Inches(0.3),
                 '(paste your dashboard screenshot here)',
                 font_size=9, color=RGBColor(156, 163, 175),
                 align=PP_ALIGN.CENTER)

    # ── RIGHT: Standard Included + intro + bullets ──
    right_x = Inches(6.0)
    right_w = Inches(3.7)
    add_text_box(s7, right_x, Inches(1.5), right_w, Inches(0.4),
                 strip_emoji(d7.get('subheadline', 'Standard Included')),
                 font_size=14, bold=True, color=BLACK, font_name=FONT_HEAD)

    intro_text = strip_emoji(d7.get('intro', '')) or 'Every conversation becomes actionable business intelligence.'
    add_text_box(s7, right_x, Inches(1.95), right_w, Inches(0.9),
                 intro_text, font_size=10, color=GRAY)

    # Bulleted metric list with C4 dots (accent 3)
    metrics = d7.get('metrics', [])
    my = 2.85
    for st in metrics[:5]:
        label = strip_emoji(st.get('label', st.get('l', '')))
        # Bullet dot (C4)
        add_oval(s7, Inches(right_x.inches), Inches(my + 0.06), Inches(0.12), Inches(0.12), C4)
        add_text_box(s7, Inches(right_x.inches + 0.22), Inches(my), Inches(right_w.inches - 0.22), Inches(0.28),
                     label, font_size=10, bold=True, color=BLACK)
        my += 0.35

    # Badges at bottom (repeated)
    if badges:
        add_text_box(s7, Inches(0.55), Inches(5.05), Inches(9.0), Inches(0.3),
                     ' | '.join([strip_emoji(b) for b in badges]),
                     font_size=9, color=GRAY, align=PP_ALIGN.CENTER)

    add_footer(s7)

    # ══════ SLIDE 15: PRICING ══════
    # Layout-15 reference:
    #   - "Pricing / Monthly" title
    #   - 3 DARK pricing cards (Starter / Premium / Enterprise)
    #   - Each tier has name, €price, "/month (12-month contract)", users + conversations,
    #     checklist of 6 features
    #   - "Addons" row under the tiers with 4 dark cards
    #
    # IMPORTANT: notso.ai pricing is FIXED. The Claude payload cannot override
    # prices / tiers / features / add-ons — only headline & lead are overridable.
    DARK_CARD = RGBColor(31, 41, 55)
    FIXED_PRICING_TIERS = [
        {
            'name': 'Starter',
            'price': '€399,-',
            'users': '1,000 unique users',
            'journeys': '3,000 conversations',
            'features': [
                '1 Journey (Customer, Client, or Event)',
                'Basic custom character design',
                'Knowledge Base: up to 3 pages',
                'Media Pack: 10 images + 5 videos',
                'Basic Analytics Portal',
                'Email support (24-hour response)',
            ],
        },
        {
            'name': 'Premium',
            'price': '€699,-',
            'users': '2,000 unique users',
            'journeys': '6,000 conversations',
            'features': [
                '2 Journeys (Customer, Client, or Event)',
                'Premium character design',
                'Knowledge Base: up to 10 pages',
                'Media Pack: 25 images + 10 videos',
                'Premium Analytics Portal',
                'Same-day email & phone support',
            ],
        },
        {
            'name': 'Enterprise',
            'price': 'Custom',
            'users': 'Unlimited users',
            'journeys': 'Unlimited conversations',
            'features': [
                'Unlimited journeys',
                'Tailor-made character design',
                'API Access / Integrations',
                'Unlimited media & branding content',
                'Custom Analytics Portal',
                '24/7 support + dedicated account manager',
            ],
        },
    ]
    FIXED_PRICING_ADDONS = [
        {'name': 'Extra Character Design', 'price': '+ €142 / month',
         'desc': 'Add one additional character design to your tier.'},
        {'name': 'Extra Journey Slot', 'price': '+ €96 / month',
         'desc': 'Add one additional customer/client/event journey.'},
        {'name': 'Partner License', 'price': '+ €149 / month per licence',
         'desc': 'Sell the service under your brand ("Notso Powered").'},
        {'name': 'Whitelabel License', 'price': '+ €349 / month per licence',
         'desc': 'Fully rebrand the platform as your own (no Notso branding).'},
    ]

    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s8, WHITE)
    d8 = proposal.get('s15', {})

    # Title + lead (headline/lead overridable, everything else locked)
    add_text_box(s8, Inches(0.55), Inches(0.3), Inches(9.0), Inches(0.55),
                 strip_emoji(d8.get('headline', 'Pricing')) + '  /  Monthly',
                 font_size=24, bold=True, color=C1, font_name=FONT_HEAD)
    _lead8 = strip_emoji(d8.get('lead', '') or d8.get('reasoning', ''))
    if _lead8:
        add_text_box(s8, Inches(0.55), Inches(0.82), Inches(9.0), Inches(0.4),
                     _lead8, font_size=10, italic=True, color=GRAY, font_name=FONT_BODY)

    tiers = FIXED_PRICING_TIERS
    num_tiers = len(tiers)
    tier_gap = 0.3
    tier_w = (8.9 - tier_gap * (num_tiers - 1)) / num_tiers
    tier_h = 3.05
    tier_top = 1.3

    # One vivid color per tier (C1 navy, C3 yellow, C4 accent)
    price_colors = [C1, C3, C4]

    for i, t in enumerate(tiers):
        tx = 0.55 + i * (tier_w + tier_gap)

        # Dark rounded card
        card = add_rounded_rect(s8, Inches(tx), Inches(tier_top),
                                Inches(tier_w), Inches(tier_h),
                                DARK_CARD, corner_radius=Inches(0.14))
        card.line.fill.background()

        # Tier name (white)
        add_text_box(s8, Inches(tx + 0.28), Inches(tier_top + 0.2),
                     Inches(tier_w - 0.56), Inches(0.35),
                     t['name'],
                     font_size=15, bold=True, color=WHITE, font_name=FONT_HEAD)

        # Price (per-tier bright color)
        pc = price_colors[i % len(price_colors)]
        price_str = t['price'].replace(' ', '\u00A0')
        add_text_box(s8, Inches(tx + 0.28), Inches(tier_top + 0.6),
                     Inches(tier_w - 0.56), Inches(0.6),
                     price_str, font_size=26, bold=True, color=pc, font_name=FONT_HEAD)

        # Period text + thin divider
        add_text_box(s8, Inches(tx + 0.28), Inches(tier_top + 1.15),
                     Inches(tier_w - 0.56), Inches(0.22),
                     '/month (12-month contract)', font_size=8,
                     color=RGBColor(156, 163, 175))
        add_rect(s8, Inches(tx + 0.28), Inches(tier_top + 1.42),
                 Inches(tier_w - 0.56), Inches(0.012),
                 RGBColor(75, 85, 99))

        # Users + conversations (stacked, bold white)
        add_text_box(s8, Inches(tx + 0.28), Inches(tier_top + 1.5),
                     Inches(tier_w - 0.56), Inches(0.2),
                     t['users'], font_size=8, bold=True, color=WHITE)
        add_text_box(s8, Inches(tx + 0.28), Inches(tier_top + 1.68),
                     Inches(tier_w - 0.56), Inches(0.2),
                     t['journeys'], font_size=8, bold=True,
                     color=RGBColor(156, 163, 175))

        # Feature checklist (white text, brand-color check)
        fy = tier_top + 1.95
        for fi, feat in enumerate(t['features']):
            add_text_box(s8, Inches(tx + 0.28), Inches(fy + fi * 0.18),
                         Inches(tier_w - 0.56), Inches(0.2),
                         '✓  ' + feat, font_size=7, color=WHITE)

    # ── Add-ons row (4 dark cards) ──
    addons = FIXED_PRICING_ADDONS
    addon_top = tier_top + tier_h + 0.18
    add_text_box(s8, Inches(0.55), Inches(addon_top), Inches(3), Inches(0.28),
                 'Addons', font_size=12, bold=True, color=C1, font_name=FONT_HEAD)
    addon_top += 0.3
    num_addons = len(addons)
    addon_gap = 0.2
    addon_w = (8.9 - addon_gap * (num_addons - 1)) / num_addons
    addon_h = 1.0

    for ai, addon in enumerate(addons):
        ax = 0.55 + ai * (addon_w + addon_gap)
        acard = add_rounded_rect(s8, Inches(ax), Inches(addon_top),
                                 Inches(addon_w), Inches(addon_h),
                                 DARK_CARD, corner_radius=Inches(0.1))
        acard.line.fill.background()
        add_text_box(s8, Inches(ax + 0.15), Inches(addon_top + 0.08),
                     Inches(addon_w - 0.3), Inches(0.22),
                     addon['name'],
                     font_size=8, bold=True, color=WHITE)
        add_text_box(s8, Inches(ax + 0.15), Inches(addon_top + 0.3),
                     Inches(addon_w - 0.3), Inches(0.22),
                     addon['price'],
                     font_size=7, bold=True, color=C3)
        add_text_box(s8, Inches(ax + 0.15), Inches(addon_top + 0.54),
                     Inches(addon_w - 0.3), Inches(0.44),
                     addon['desc'],
                     font_size=6, color=RGBColor(156, 163, 175))

    add_footer(s8)

    # ══════ SLIDE 18: THANK YOU ══════
    # Notso reference style (Fin slide 12):
    #   - Black background
    #   - Huge "Bedankt!" / "Thank You!" title (80pt+)
    #   - Closing message paragraph
    #   - Contact info: phone, email, website (large, bold)
    #   - Company details (small, gray)
    #   - Mascot image on the right
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    DARK_BG = RGBColor(0, 0, 0)
    DARK_TEXT = RGBColor(30, 30, 30)  # near-black text on colored card
    add_bg(s9, DARK_BG)
    d9 = proposal.get('s18', {})

    # Large inner card with brand color (spanning most of the slide)
    card_margin_x = Inches(0.35)
    card_margin_y = Inches(0.35)
    card_w = W - card_margin_x * 2
    card_h = H - card_margin_y * 2
    add_rounded_rect(s9, card_margin_x, card_margin_y, card_w, card_h, C1)

    # LEFT COLUMN (0.8" to 5.8") — text only, no overlap with mascot
    # RIGHT COLUMN (6.2" to 9.5") — mascot image + company details

    # ── Huge title: "Thank You!" ──
    _closing_title = strip_emoji(d9.get('closing_title', 'Thank You!'))
    add_text_box(s9, Inches(0.8), Inches(0.6), Inches(4.8), Inches(1.2),
                 _closing_title, font_size=56, bold=False, color=DARK_TEXT, font_name=FONT_HEAD)

    # ── Closing message ──
    _lead9 = strip_emoji(d9.get('lead', ''))
    closing_text = strip_emoji(d9.get('closing', ''))
    combined_closing = (_lead9 + '\n' + closing_text).strip() if _lead9 else closing_text
    # Truncate if too long to prevent overflow
    if len(combined_closing) > 200:
        combined_closing = combined_closing[:197] + '...'
    if combined_closing:
        add_text_box(s9, Inches(0.8), Inches(1.75), Inches(4.8), Inches(1.3),
                     combined_closing, font_size=10, color=DARK_TEXT, font_name=FONT_BODY)

    # ── Contact info (large, bold) ──
    info_items = [
        d9.get('phone', '+31 6 34 197 668'),
        d9.get('email', 'hello@notso.ai'),
        d9.get('website', 'www.notso.ai'),
    ]
    iy = 3.15
    for val in info_items:
        if val:
            add_text_box(s9, Inches(0.8), Inches(iy), Inches(4.8), Inches(0.4),
                         strip_emoji(val), font_size=16, bold=True, color=DARK_TEXT, font_name=FONT_HEAD)
            iy += 0.42

    # ── Mascot image (right column) ──
    add_placeholder_img(s9, Inches(6.2), Inches(0.5), Inches(3.2), Inches(3.2),
                        image_key='cover', label='Mascot')

    # ── Company details (right column, below mascot) ──
    add_text_box(s9, Inches(6.2), Inches(3.9), Inches(1.5), Inches(1.0),
                 'Notso B.V.\nBerkenstraat 11\nDuivendrecht',
                 font_size=7, color=DARK_TEXT)
    add_text_box(s9, Inches(7.8), Inches(3.9), Inches(1.6), Inches(1.0),
                 'KvK: 95131124\nRSIN: 867013916\nVAT: NL867013316B01',
                 font_size=7, color=DARK_TEXT)

    # ── Next steps (bottom, if any) ──
    steps = d9.get('next_steps', [])
    if steps:
        steps_text = '  |  '.join([strip_emoji(s) for s in steps[:3]])
        add_text_box(s9, Inches(0.8), Inches(5.0), Inches(8.5), Inches(0.4),
                     'Next: ' + steps_text, font_size=9, color=DARK_TEXT,
                     align=PP_ALIGN.LEFT)

    # ══════ NEW SLIDES (s2, s4, s8, s10, s11, s13, s14, s16, s17) ══════
    # These are added after the existing 9 slides, then we reorder and filter

    # s2: Table of Contents
    # Notso reference style: clean light background, black text, brand-colored numbers
    if 's2' in sel:
        stoc = prs.slides.add_slide(prs.slide_layouts[6])
        TOC_BG = RGBColor(244, 244, 243)
        add_bg(stoc, TOC_BG)

        # Huge left-aligned headline (matching Fin style: 45pt black Poppins)
        add_text_box(stoc, Inches(0.6), Inches(0.35), Inches(5.0), Inches(0.8),
                     'TABLE OF', font_size=40, bold=True, color=BLACK,
                     align=PP_ALIGN.LEFT, font_name=FONT_HEAD)
        add_text_box(stoc, Inches(0.6), Inches(1.0), Inches(5.0), Inches(0.8),
                     'CONTENTS', font_size=40, bold=True, color=BLACK,
                     align=PP_ALIGN.LEFT, font_name=FONT_HEAD)

        # TOC = all selected slides except Cover, TOC itself, Thank You → 15 items
        exclude_ids = {'s1', 's2', 's18'}
        toc_items = [s for s in ['s3','s4','s5','s6','s7','s8','s9','s10','s11','s12','s13','s14','s15','s16','s17']
                     if s in sel and s not in exclude_ids]
        toc_labels = {'s3':'Pain Points','s4':'Market Opportunity','s5':'Core Features',
                      's6':'Mascot Selection','s7':'Mascot Design','s8':'Personality & Empathy',
                      's9':'Chat Mock-up','s10':'Chatflow Design','s11':'Knowledge Base',
                      's12':'Data & Insights','s13':'ROI Evidence','s14':'Roadmap',
                      's15':'Pricing','s16':'Promo Materials','s17':'Licensing'}

        # 4-column grid, filled TOP-TO-BOTTOM (column-major) so readers scan down first
        cols = 4
        col_w = 2.15
        row_h = 0.72
        grid_top = 2.25
        grid_left = 0.55
        n_items = len(toc_items)
        # rows per column, rounding up
        rows_per_col = (n_items + cols - 1) // cols if n_items else 1
        for i, s in enumerate(toc_items):
            col = i // rows_per_col
            row = i % rows_per_col
            x = grid_left + col * col_w
            y = grid_top + row * row_h
            add_text_box(stoc, Inches(x), Inches(y), Inches(0.55), Inches(0.35),
                         f'{i+1:02d}', font_size=12, bold=True, color=C3, font_name=FONT_HEAD)
            add_text_box(stoc, Inches(x), Inches(y + 0.3), Inches(col_w - 0.1), Inches(0.4),
                         toc_labels.get(s, s).upper(),
                         font_size=11, bold=True, color=BLACK, font_name=FONT_HEAD)

    # s4: Market Opportunity
    # Layout-04 reference: headline + lead, 3 stat cards with small "Industry Size"
    # labels in C3 (yellow) and huge numbers in C1, sources in gray, then a
    # rounded "Competitive Gaps" card below.
    if 's4' in sel:
        d4m = proposal.get('s4', {})
        sm = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(sm, LIGHT_BG)

        after_title = add_hero_band(sm,
            d4m.get('headline', 'Market Opportunity'),
            d4m.get('lead', '') or d4m.get('opportunity', ''),
            band_color=C1, title_size=22, band_height=1.25)

        mdata = [('Industry Size', d4m.get('industry_size',{})),
                 ('Growth Rate', d4m.get('growth_rate',{})),
                 ('Projected Size', d4m.get('projected_size',{}))]
        stat_y = max(1.8, after_title + 0.25)
        for i, (label, data) in enumerate(mdata):
            if data:
                x = 0.55 + i * 3.1
                # Label as a C3 YELLOW PILL with BLACK text (not yellow text)
                label_w = min(1.65, 0.3 + len(label) * 0.11)
                pill = add_rounded_rect(sm, Inches(x), Inches(stat_y),
                                        Inches(label_w), Inches(0.32),
                                        C3, corner_radius=Inches(0.04))
                pill.line.fill.background()
                add_text_box(sm, Inches(x), Inches(stat_y + 0.03), Inches(label_w), Inches(0.28),
                             label, font_size=10, bold=True, color=BLACK,
                             align=PP_ALIGN.CENTER, font_name=FONT_HEAD)
                _val = strip_emoji(str(data.get('value','')))
                _fs = 32 if len(_val) <= 8 else (26 if len(_val) <= 12 else 20)
                # Up-arrow prefix (references layout with ↑ prefix)
                display = '↑ ' + _val
                add_text_box(sm, Inches(x), Inches(stat_y + 0.45), Inches(2.85), Inches(0.85),
                             display, font_size=_fs, bold=True, color=BLACK, font_name=FONT_HEAD)
                add_text_box(sm, Inches(x), Inches(stat_y + 1.32), Inches(2.85), Inches(0.3),
                             strip_emoji(str(data.get('label', label))),
                             font_size=9, color=BLACK)
                add_text_box(sm, Inches(x), Inches(stat_y + 1.55), Inches(2.85), Inches(0.3),
                             'Source: ' + strip_emoji(str(data.get('source',''))),
                             font_size=8, color=GRAY)

        # Competitive Gaps card (rounded, white on the light-gray background)
        comps = d4m.get('competitors', [])
        if comps:
            gap_y = stat_y + 2.05
            gap_card = add_rounded_rect(sm, Inches(0.55), Inches(gap_y),
                                        Inches(8.9), Inches(1.25),
                                        WHITE, corner_radius=Inches(0.2))
            gap_card.line.color.rgb = BORDER
            gap_card.line.width = Pt(0.75)
            add_text_box(sm, Inches(0.8), Inches(gap_y + 0.1), Inches(8.4), Inches(0.3),
                         'Competitive Gaps', font_size=12, bold=True, color=BLACK, font_name=FONT_HEAD)
            comp_text = '\n'.join([
                f"• {strip_emoji(c.get('name',''))}: {strip_emoji(c.get('gap',''))}"
                for c in comps[:3]
            ])
            add_text_box(sm, Inches(0.8), Inches(gap_y + 0.4), Inches(8.4), Inches(0.9),
                         comp_text, font_size=9, color=GRAY)
        add_footer(sm)

    # s8: Personality & Empathy (Expression Grid)
    # Layout-08 reference: left column has a text intro (name + summary),
    # right side has a 3x3 grid of pastel-tinted cards (each a different color).
    if 's8' in sel:
        d8e = proposal.get('s8', {})
        se = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(se, WHITE)
        _name8 = strip_emoji(d8e.get('name', '')) or 'mascot'
        add_title_lead(se,
            'Personality & Empathy — ' + _name8,
            d8e.get('lead', '') or d8e.get('personality_summary', ''),
            top=0.35, title_size=22)

        # Cycle through tinted variants of all 4 brand colors for pastel cards
        palette = [C1_TINT, C3_TINT, C2_TINT, C3_TINT, C1_TINT, C4_TINT,
                   C2_TINT, C4_TINT, C1_TINT]
        expressions = d8e.get('expressions', [])
        grid_left = 2.8
        cell_w = 2.15
        cell_h = 1.15
        cell_gap = 0.15
        grid_top = 1.6
        for i, expr in enumerate(expressions[:9]):
            row, col = divmod(i, 3)
            x = grid_left + col * (cell_w + cell_gap)
            y = grid_top + row * (cell_h + cell_gap)
            exp_card = add_rect(se, Inches(x), Inches(y), Inches(cell_w), Inches(cell_h),
                                palette[i % len(palette)])
            exp_card.line.fill.background()
            add_text_box(se, Inches(x + 0.15), Inches(y + 0.1), Inches(cell_w - 0.3), Inches(0.35),
                         strip_emoji(expr.get('emotion','')),
                         font_size=12, bold=True, color=BLACK, font_name=FONT_HEAD)
            add_text_box(se, Inches(x + 0.15), Inches(y + 0.45), Inches(cell_w - 0.3), Inches(cell_h - 0.55),
                         strip_emoji(expr.get('context','')),
                         font_size=8, color=BLACK)

        # Left panel — bottom intro paragraph
        add_text_box(se, Inches(0.55), Inches(4.0), Inches(2.15), Inches(1.4),
                     strip_emoji(d8e.get('personality_summary','')),
                     font_size=9, color=BLACK)
        add_footer(se)

    # s10: Chatflow Design
    # Layout-10 reference: hero band with title on blue, then a horizontal flow of
    # dark rounded "node cards" with colored icon pills on top, connecting chevron
    # lines between each pair. Each card has an icon label, step title (white),
    # and a short description (light gray).
    if 's10' in sel:
        d10 = proposal.get('s10', {})
        scf = prs.slides.add_slide(prs.slide_layouts[6])

        # Dark theme colors
        DARK_BG = RGBColor(8, 11, 18)
        DARK_NODE = RGBColor(24, 30, 42)
        DARKER = RGBColor(32, 40, 54)
        BORDER_DARK = RGBColor(75, 85, 99)
        SOFT_WHITE = RGBColor(229, 231, 235)
        MUTED = RGBColor(156, 163, 175)
        HIGHLIGHT_GREEN = RGBColor(34, 197, 94)
        BADGE_PURPLE = RGBColor(124, 58, 237)

        add_bg(scf, DARK_BG)

        # Title + lead (white on dark)
        add_text_box(scf, Inches(0.55), Inches(0.28), Inches(9), Inches(0.55),
                     strip_emoji(d10.get('headline', 'How the conversation flows')),
                     font_size=22, bold=True, color=WHITE, font_name=FONT_HEAD)
        lead = strip_emoji(d10.get('lead', ''))
        if lead:
            add_text_box(scf, Inches(0.55), Inches(0.82), Inches(9), Inches(0.35),
                         lead, font_size=11, italic=True, color=SOFT_WHITE,
                         font_name=FONT_BODY)

        # ── Flow column data ──
        # New branching format: d10.flow.columns[] where each column is a single
        # node or a stack of 2 nodes. Each node has a `type`:
        #   bot_msg / user_msg / options / badge / link
        # Falls back to legacy d10.steps if no flow provided.
        flow = d10.get('flow', {}) or {}
        columns = flow.get('columns', [])
        if not columns:
            steps = d10.get('steps', [])
            fallback_types = ['bot_msg', 'user_msg', 'options', 'badge', 'bot_msg']
            columns = []
            for i, s in enumerate(steps[:5]):
                columns.append({
                    'type': fallback_types[i % len(fallback_types)],
                    'title': s.get('title', ''),
                    'text': s.get('desc', ''),
                })

        n_cols = max(len(columns), 1)

        # Column geometry
        flow_top = 1.35
        flow_bottom = 5.3
        flow_h = flow_bottom - flow_top
        flow_mid = flow_top + flow_h / 2
        x_start = 0.35
        x_end = 9.65
        col_gap = 0.12
        col_w = min(1.55, (x_end - x_start - col_gap * (n_cols - 1)) / n_cols)
        total_flow_w = n_cols * col_w + (n_cols - 1) * col_gap
        flow_x0 = (10 - total_flow_w) / 2

        node_color_cycle = [C1, C3, C4, C2]

        def col_x(i):
            return flow_x0 + i * (col_w + col_gap)

        def draw_avatar(x, y, size, color):
            av = add_rounded_rect(scf, Inches(x), Inches(y),
                                  Inches(size), Inches(size),
                                  color, corner_radius=Inches(size / 2))
            av.line.fill.background()

        def draw_bot_msg(i, node, y_top, max_h):
            """Dark card with colored top bar + title, speech bubble below."""
            x = col_x(i)
            card = add_rounded_rect(scf, Inches(x), Inches(y_top),
                                    Inches(col_w), Inches(max_h),
                                    DARK_NODE, corner_radius=Inches(0.12))
            card.line.color.rgb = BORDER_DARK
            card.line.width = Pt(0.75)
            # Small accent square top-left
            pill_color = node_color_cycle[i % 4]
            accent = add_rounded_rect(scf, Inches(x + 0.12), Inches(y_top + 0.12),
                                       Inches(0.18), Inches(0.18), pill_color,
                                       corner_radius=Inches(0.04))
            accent.line.fill.background()
            # Title on its own row (full card width) so long titles don't break
            add_text_box(scf, Inches(x + 0.12), Inches(y_top + 0.32),
                         Inches(col_w - 0.24), Inches(0.26),
                         strip_emoji(node.get('title', '')),
                         font_size=9, bold=True, color=WHITE, font_name=FONT_HEAD)
            # Speech bubble (darker rounded box)
            bubble_top = y_top + 0.62
            tags = node.get('tags', []) or []
            bubble_h_reserve = 0.25 * min(len(tags), 2) + (0.08 if tags else 0)
            bubble_h = max_h - (bubble_top - y_top) - 0.12 - bubble_h_reserve
            bubble = add_rounded_rect(scf, Inches(x + 0.1), Inches(bubble_top),
                                      Inches(col_w - 0.2), Inches(bubble_h),
                                      DARKER, corner_radius=Inches(0.08))
            bubble.line.color.rgb = BORDER_DARK
            bubble.line.width = Pt(0.5)
            add_text_box(scf, Inches(x + 0.16), Inches(bubble_top + 0.06),
                         Inches(col_w - 0.32), Inches(bubble_h - 0.12),
                         strip_emoji(node.get('text', '')),
                         font_size=8, color=SOFT_WHITE)
            # Optional tag pills below the bubble
            if tags:
                tag_start_y = bubble_top + bubble_h + 0.06
                for ti, tag in enumerate(tags[:2]):
                    ty = tag_start_y + ti * 0.22
                    pill = add_rounded_rect(scf, Inches(x + 0.12), Inches(ty),
                                            Inches(col_w - 0.24), Inches(0.18),
                                            DARK_BG, corner_radius=Inches(0.09))
                    pill.line.color.rgb = BADGE_PURPLE
                    pill.line.width = Pt(0.75)
                    add_text_box(scf, Inches(x + 0.12), Inches(ty + 0.01),
                                 Inches(col_w - 0.24), Inches(0.16),
                                 strip_emoji(tag), font_size=7, bold=True,
                                 color=BADGE_PURPLE, align=PP_ALIGN.CENTER,
                                 font_name=FONT_HEAD)

        def draw_user_msg(i, node, y_top, max_h):
            """User chat bubble (white message bubble on dark card)."""
            x = col_x(i)
            card = add_rounded_rect(scf, Inches(x), Inches(y_top),
                                    Inches(col_w), Inches(max_h),
                                    DARK_NODE, corner_radius=Inches(0.12))
            card.line.color.rgb = BORDER_DARK
            card.line.width = Pt(0.75)
            # Gray accent square top-left
            accent = add_rounded_rect(scf, Inches(x + 0.12), Inches(y_top + 0.12),
                                       Inches(0.18), Inches(0.18),
                                       RGBColor(107, 114, 128),
                                       corner_radius=Inches(0.04))
            accent.line.fill.background()
            add_text_box(scf, Inches(x + 0.12), Inches(y_top + 0.32),
                         Inches(col_w - 0.24), Inches(0.26),
                         strip_emoji(node.get('title', 'User')),
                         font_size=9, bold=True, color=WHITE, font_name=FONT_HEAD)
            # White message bubble
            bubble_top = y_top + 0.62
            bubble_h = max_h - 0.74
            bubble = add_rounded_rect(scf, Inches(x + 0.1), Inches(bubble_top),
                                      Inches(col_w - 0.2), Inches(bubble_h),
                                      WHITE, corner_radius=Inches(0.08))
            bubble.line.fill.background()
            add_text_box(scf, Inches(x + 0.16), Inches(bubble_top + 0.06),
                         Inches(col_w - 0.32), Inches(bubble_h - 0.12),
                         strip_emoji(node.get('text', '')),
                         font_size=8, color=BLACK)

        def draw_options(i, node, y_top, max_h):
            """Box with title + stacked option pills; one can be highlighted."""
            x = col_x(i)
            card = add_rounded_rect(scf, Inches(x), Inches(y_top),
                                    Inches(col_w), Inches(max_h),
                                    DARK_NODE, corner_radius=Inches(0.12))
            card.line.color.rgb = BORDER_DARK
            card.line.width = Pt(0.75)
            add_text_box(scf, Inches(x + 0.1), Inches(y_top + 0.1),
                         Inches(col_w - 0.2), Inches(0.32),
                         strip_emoji(node.get('title', '')),
                         font_size=10, bold=True, color=WHITE, font_name=FONT_HEAD,
                         align=PP_ALIGN.CENTER)
            items = node.get('items', [])
            highlight_idx = node.get('highlight', -1)
            opt_top = y_top + 0.5
            opt_avail = max_h - 0.6
            opt_h = min(0.28, (opt_avail - 0.04 * max(len(items) - 1, 0)) / max(len(items), 1))
            for oi, it in enumerate(items[:5]):
                oy = opt_top + oi * (opt_h + 0.05)
                is_hl = (oi == highlight_idx)
                pill_bg = DARK_BG
                pill = add_rounded_rect(scf, Inches(x + 0.15), Inches(oy),
                                        Inches(col_w - 0.3), Inches(opt_h),
                                        pill_bg, corner_radius=Inches(opt_h / 2))
                pill.line.color.rgb = HIGHLIGHT_GREEN if is_hl else BORDER_DARK
                pill.line.width = Pt(1.25 if is_hl else 0.75)
                add_text_box(scf, Inches(x + 0.15), Inches(oy + 0.04),
                             Inches(col_w - 0.3), Inches(opt_h - 0.08),
                             strip_emoji(it), font_size=8,
                             bold=is_hl,
                             color=HIGHLIGHT_GREEN if is_hl else SOFT_WHITE,
                             align=PP_ALIGN.CENTER)

        def draw_badge(i, node, y_top, max_h):
            """Compact centered badge (like Control Agent)."""
            x = col_x(i)
            badge_h = 0.7
            badge_y = y_top + max_h / 2 - badge_h / 2
            card = add_rounded_rect(scf, Inches(x + 0.05), Inches(badge_y),
                                    Inches(col_w - 0.1), Inches(badge_h),
                                    DARK_NODE, corner_radius=Inches(0.1))
            card.line.color.rgb = C3
            card.line.width = Pt(1.5)
            # Check icon
            add_text_box(scf, Inches(x + 0.05), Inches(badge_y + 0.06),
                         Inches(col_w - 0.1), Inches(0.3),
                         '✓', font_size=14, bold=True, color=C3,
                         align=PP_ALIGN.CENTER)
            add_text_box(scf, Inches(x + 0.05), Inches(badge_y + 0.35),
                         Inches(col_w - 0.1), Inches(0.3),
                         strip_emoji(node.get('title', '')),
                         font_size=9, bold=True, color=C3,
                         align=PP_ALIGN.CENTER, font_name=FONT_HEAD)

        def draw_link(i, node, y_top, max_h):
            """Link-style button (bright blue border)."""
            x = col_x(i)
            btn_h = 0.55
            btn_y = y_top + max_h / 2 - btn_h / 2
            card = add_rounded_rect(scf, Inches(x + 0.05), Inches(btn_y),
                                    Inches(col_w - 0.1), Inches(btn_h),
                                    DARK_NODE, corner_radius=Inches(0.1))
            card.line.color.rgb = RGBColor(59, 130, 246)
            card.line.width = Pt(1.25)
            add_text_box(scf, Inches(x + 0.05), Inches(btn_y + 0.16),
                         Inches(col_w - 0.1), Inches(0.28),
                         '🔗  ' + strip_emoji(node.get('title', 'Link to module')),
                         font_size=9, bold=True, color=RGBColor(96, 165, 250),
                         align=PP_ALIGN.CENTER, font_name=FONT_HEAD)

        def draw_node(i, node, y_top, max_h):
            t = node.get('type', 'bot_msg')
            if t == 'user_msg':
                draw_user_msg(i, node, y_top, max_h)
            elif t == 'options':
                draw_options(i, node, y_top, max_h)
            elif t == 'badge':
                draw_badge(i, node, y_top, max_h)
            elif t == 'link':
                draw_link(i, node, y_top, max_h)
            else:
                draw_bot_msg(i, node, y_top, max_h)

        # Layout each column: stack mode supports 2 nodes in one column
        col_centers = []  # list of (x_right_edge, y_mid_list) for connecting lines
        for i, col in enumerate(columns):
            x = col_x(i)
            if col.get('type') == 'stack':
                sub_nodes = col.get('nodes', [])[:2]
                n_sub = max(len(sub_nodes), 1)
                sub_gap = 0.15
                sub_h = (flow_h - sub_gap * (n_sub - 1)) / n_sub
                mids = []
                for si, sn in enumerate(sub_nodes):
                    sy = flow_top + si * (sub_h + sub_gap)
                    draw_node(i, sn, sy, sub_h)
                    mids.append(sy + sub_h / 2)
                col_centers.append((x, x + col_w, mids))
            else:
                draw_node(i, col, flow_top, flow_h)
                col_centers.append((x, x + col_w, [flow_mid]))

        # ── Draw connector lines between columns ──
        # Each connection is a thin horizontal/L-shape from right edge of col[i]
        # to left edge of col[i+1]. If one side has multiple mids, branch.
        for i in range(n_cols - 1):
            (x0, right_edge_a, mids_a) = col_centers[i]
            (left_edge_b, x1, mids_b) = col_centers[i + 1]
            # Horizontal span between the two cards
            span_x = right_edge_a
            span_w = left_edge_b - right_edge_a
            # Branch out from single to multi, or converge from multi to single
            from_mids = mids_a
            to_mids = mids_b
            # Connect each "from" to each "to" (if both are singles, just 1 line)
            for fm in from_mids:
                for tm in to_mids:
                    # Horizontal line at average y (or L-shape if different)
                    if abs(fm - tm) < 0.02:
                        add_rect(scf, Inches(span_x), Inches(fm - 0.008),
                                 Inches(span_w), Inches(0.016), MUTED)
                    else:
                        # L-shape: half horizontal, vertical, half horizontal
                        mid_x = span_x + span_w / 2
                        # First half horizontal at fm
                        add_rect(scf, Inches(span_x), Inches(fm - 0.008),
                                 Inches(span_w / 2 + 0.008), Inches(0.016), MUTED)
                        # Vertical segment
                        y_lo = min(fm, tm)
                        y_hi = max(fm, tm)
                        add_rect(scf, Inches(mid_x - 0.008), Inches(y_lo),
                                 Inches(0.016), Inches(y_hi - y_lo), MUTED)
                        # Second half horizontal at tm
                        add_rect(scf, Inches(mid_x - 0.008), Inches(tm - 0.008),
                                 Inches(span_w / 2 + 0.016), Inches(0.016), MUTED)
        add_footer(scf)

    # s11: Knowledge Base
    # Layout-11 reference: left column "提供文件" input cards (yellow/beige headers)
    # -> middle mascot slot with arrow -> right column extracted knowledge cards
    # (navy headers). 3 rows of input→output pairs.
    if 's11' in sel:
        d11 = proposal.get('s11', {})
        skb = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(skb, LIGHT_BG)
        after_title = add_hero_band(skb,
            d11.get('headline', 'Knowledge Base'),
            d11.get('lead', ''),
            band_color=C1, title_size=22, band_height=1.25)

        cats = d11.get('categories', [])[:3]
        n_cats = max(len(cats), 1)

        # Column geometry
        left_col_x = 0.55
        left_col_w = 3.55
        mid_col_x = 4.25
        mid_col_w = 1.5
        right_col_x = 5.9
        right_col_w = 3.55

        # Column subheaders
        sub_y = after_title + 0.15
        add_text_box(skb, Inches(left_col_x), Inches(sub_y),
                     Inches(left_col_w), Inches(0.3),
                     'Input · Documents',
                     font_size=10, bold=True, color=GRAY, font_name=FONT_HEAD)
        add_text_box(skb, Inches(right_col_x), Inches(sub_y),
                     Inches(right_col_w), Inches(0.3),
                     'Output · Extracted knowledge',
                     font_size=10, bold=True, color=C1, font_name=FONT_HEAD,
                     align=PP_ALIGN.RIGHT)

        rows_top = sub_y + 0.45
        row_gap = 0.18
        available_h = 5.25 - rows_top
        row_h = (available_h - row_gap * (n_cats - 1)) / n_cats
        row_h = max(0.9, min(row_h, 1.25))

        for i, cat in enumerate(cats):
            y = rows_top + i * (row_h + row_gap)

            # LEFT input card (yellow/beige header, white body)
            lcard = add_rounded_rect(skb, Inches(left_col_x), Inches(y),
                                     Inches(left_col_w), Inches(row_h),
                                     WHITE, corner_radius=Inches(0.1))
            lcard.line.color.rgb = BORDER
            lcard.line.width = Pt(0.75)
            # Yellow (C3) left vertical accent bar
            add_rect(skb, Inches(left_col_x), Inches(y),
                     Inches(0.1), Inches(row_h), C3)
            # Header title (yellow pill-like tag)
            tag_w = 1.05
            add_rounded_rect(skb, Inches(left_col_x + 0.22), Inches(y + 0.15),
                             Inches(tag_w), Inches(0.28), C3,
                             corner_radius=Inches(0.05))
            add_text_box(skb, Inches(left_col_x + 0.22), Inches(y + 0.17),
                         Inches(tag_w), Inches(0.25),
                         'Documents', font_size=9, bold=True, color=BLACK,
                         align=PP_ALIGN.CENTER, font_name=FONT_HEAD)
            add_text_box(skb, Inches(left_col_x + 1.35), Inches(y + 0.14),
                         Inches(left_col_w - 1.5), Inches(0.32),
                         strip_emoji(cat.get('title', '')),
                         font_size=11, bold=True, color=BLACK, font_name=FONT_HEAD)
            # Doc list
            docs = cat.get('docs', []) or []
            if not docs:
                docs = cat.get('items', [])[:2]
            docs_text = '\n'.join(['📄  ' + strip_emoji(d) for d in docs[:4]])
            add_text_box(skb, Inches(left_col_x + 0.22), Inches(y + 0.5),
                         Inches(left_col_w - 0.4), Inches(row_h - 0.58),
                         docs_text, font_size=9, color=GRAY)

            # MIDDLE arrow connector (red) pointing right
            arrow_y = y + row_h / 2 - 0.12
            add_rect(skb, Inches(mid_col_x + 0.15), Inches(arrow_y + 0.08),
                     Inches(mid_col_w - 0.45), Inches(0.08), C4)
            add_text_box(skb, Inches(mid_col_x + mid_col_w - 0.55),
                         Inches(arrow_y - 0.05), Inches(0.5), Inches(0.35),
                         '▶', font_size=18, bold=True, color=C4,
                         align=PP_ALIGN.CENTER)

            # RIGHT output card (navy header, white body)
            rcard = add_rounded_rect(skb, Inches(right_col_x), Inches(y),
                                     Inches(right_col_w), Inches(row_h),
                                     WHITE, corner_radius=Inches(0.1))
            rcard.line.color.rgb = BORDER
            rcard.line.width = Pt(0.75)
            # Navy (C1) header bar
            add_rounded_rect(skb, Inches(right_col_x), Inches(y),
                             Inches(right_col_w), Inches(0.38), C1,
                             corner_radius=Inches(0.1))
            add_text_box(skb, Inches(right_col_x + 0.22), Inches(y + 0.05),
                         Inches(right_col_w - 0.4), Inches(0.32),
                         f'0{i+1}  ' + strip_emoji(cat.get('title', '')),
                         font_size=12, bold=True, color=WHITE, font_name=FONT_HEAD)
            # Bullet items (extracted knowledge)
            items = cat.get('items', []) or []
            items_text = '\n'.join(['•  ' + strip_emoji(it) for it in items[:4]])
            add_text_box(skb, Inches(right_col_x + 0.22), Inches(y + 0.48),
                         Inches(right_col_w - 0.4), Inches(row_h - 0.55),
                         items_text, font_size=10, color=BLACK)

        # Central mascot slot (dashed light-blue pill in the middle column)
        mascot_slot_h = 1.7
        mascot_slot_y = rows_top + (n_cats * (row_h + row_gap) - row_gap) / 2 - mascot_slot_h / 2
        mslot = add_rounded_rect(skb, Inches(mid_col_x + 0.05),
                                 Inches(mascot_slot_y),
                                 Inches(mid_col_w - 0.1),
                                 Inches(mascot_slot_h),
                                 C1_TINT, corner_radius=Inches(0.18))
        mslot.line.color.rgb = C1
        mslot.line.width = Pt(1.25)
        try:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            mslot.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass
        add_text_box(skb, Inches(mid_col_x), Inches(mascot_slot_y + mascot_slot_h / 2 - 0.18),
                     Inches(mid_col_w), Inches(0.35),
                     'Mascot',
                     font_size=10, italic=True, color=C1,
                     align=PP_ALIGN.CENTER, font_name=FONT_BODY)
        add_footer(skb)

    # s13: ROI Evidence
    # Layout-13 reference: 4 stat cards with huge C1 numbers, then a Before|After
    # comparison strip with row labels on the left (C3-style pill).
    if 's13' in sel:
        d13 = proposal.get('s13', {})
        sroi = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(sroi, WHITE)
        after_title = add_title_lead(sroi,
            d13.get('headline', 'Proven ROI'),
            d13.get('lead', ''),
            top=0.35, title_size=24)

        roi_stats = d13.get('stats', [])
        n_stats = max(len(roi_stats[:4]), 1)
        gap = 0.2
        stat_w = (8.9 - gap * (n_stats - 1)) / n_stats
        stat_top = max(1.3, after_title + 0.05)
        stat_h = 1.2
        # Cycle through tinted pastels + matching accent numbers so all 4 brand colors appear
        bg_cycle = [C1_TINT, C3_TINT, C4_TINT, C2_TINT]
        num_cycle = [C1, C3, C4, C2]
        for i, st in enumerate(roi_stats[:4]):
            x = 0.55 + i * (stat_w + gap)
            card = add_rounded_rect(sroi, Inches(x), Inches(stat_top), Inches(stat_w), Inches(stat_h),
                                    bg_cycle[i % len(bg_cycle)], corner_radius=Inches(0.1))
            card.line.fill.background()
            # Label at top (bold black)
            add_text_box(sroi, Inches(x + 0.1), Inches(stat_top + 0.08), Inches(stat_w - 0.2), Inches(0.3),
                         strip_emoji(st.get('l','')),
                         font_size=10, bold=True, color=BLACK,
                         align=PP_ALIGN.CENTER, font_name=FONT_HEAD)
            # Huge number in matching brand color
            add_text_box(sroi, Inches(x + 0.1), Inches(stat_top + 0.36), Inches(stat_w - 0.2), Inches(0.6),
                         strip_emoji(st.get('n','')),
                         font_size=26, bold=True, color=num_cycle[i % len(num_cycle)],
                         align=PP_ALIGN.CENTER, font_name=FONT_HEAD)
            # Detail line beneath number (if provided in payload)
            detail = strip_emoji(st.get('detail', ''))
            if detail:
                add_text_box(sroi, Inches(x + 0.1), Inches(stat_top + 0.92), Inches(stat_w - 0.2), Inches(0.25),
                             detail, font_size=9, italic=True, color=GRAY,
                             align=PP_ALIGN.CENTER)
        # Before / After — row-based layout (label pill | before | arrow | after)
        ba = d13.get('before_after', {})
        if ba:
            # Support both new `rows` format and legacy before[]/after[] arrays
            rows = ba.get('rows')
            if not rows:
                before_list = ba.get('before', [])
                after_list = ba.get('after', [])
                default_labels = ba.get('labels', [])
                n_rows = min(len(before_list), len(after_list))
                rows = []
                for i in range(n_rows):
                    rows.append({
                        'label': default_labels[i] if i < len(default_labels) else f'Item {i+1}',
                        'before': before_list[i],
                        'after': after_list[i],
                    })
            rows = rows[:3]

            ba_top = stat_top + stat_h + 0.18
            # Column x positions
            label_x = 0.55
            label_w = 1.35
            before_x = 2.0
            before_w = 3.05
            arrow_x = 5.1
            arrow_w = 0.5
            after_x = 5.65
            after_w = 3.8

            # Column headers "Before" and "After notso.ai"
            header_y = ba_top
            add_text_box(sroi, Inches(before_x), Inches(header_y),
                         Inches(before_w), Inches(0.3),
                         'Before', font_size=13, bold=True, color=C1,
                         align=PP_ALIGN.CENTER, font_name=FONT_HEAD)
            add_text_box(sroi, Inches(after_x), Inches(header_y),
                         Inches(after_w), Inches(0.3),
                         'After notso.ai', font_size=13, bold=True, color=C1,
                         align=PP_ALIGN.CENTER, font_name=FONT_HEAD)

            # Row area starts below the header
            rows_top = header_y + 0.32
            rows_bottom = 5.3
            avail_h = rows_bottom - rows_top
            n_r = max(len(rows), 1)
            row_h = avail_h / n_r
            row_line_w = 8.95

            # Top divider line (above first row)
            add_rect(sroi, Inches(label_x), Inches(rows_top - 0.01),
                     Inches(row_line_w), Inches(0.012), BORDER)

            for i, rw in enumerate(rows):
                y = rows_top + i * row_h
                row_mid_y = y + row_h / 2

                # Label pill on the left (navy rounded) — sized to row height
                pill_h = min(row_h - 0.16, 0.62)
                pill_y = row_mid_y - pill_h / 2
                add_rounded_rect(sroi, Inches(label_x), Inches(pill_y),
                                 Inches(label_w), Inches(pill_h), C1,
                                 corner_radius=Inches(0.08))
                add_text_box(sroi, Inches(label_x + 0.05), Inches(pill_y + 0.04),
                             Inches(label_w - 0.1), Inches(pill_h - 0.08),
                             strip_emoji(rw.get('label', '')),
                             font_size=11, bold=True, color=WHITE,
                             align=PP_ALIGN.CENTER, font_name=FONT_HEAD)

                # Before text (bulleted, gray)
                before_text = rw.get('before', '')
                if isinstance(before_text, list):
                    before_text = '\n'.join(['•  ' + strip_emoji(b) for b in before_text])
                else:
                    before_text = '•  ' + strip_emoji(before_text)
                add_text_box(sroi, Inches(before_x), Inches(y + 0.1),
                             Inches(before_w), Inches(row_h - 0.15),
                             before_text, font_size=10, color=GRAY)

                # Middle arrow →
                add_text_box(sroi, Inches(arrow_x), Inches(row_mid_y - 0.2),
                             Inches(arrow_w), Inches(0.4),
                             '→', font_size=22, bold=True, color=C1,
                             align=PP_ALIGN.CENTER, font_name=FONT_HEAD)

                # After text (bulleted, black)
                after_text = rw.get('after', '')
                if isinstance(after_text, list):
                    after_text = '\n'.join(['•  ' + strip_emoji(a) for a in after_text])
                else:
                    after_text = '•  ' + strip_emoji(after_text)
                add_text_box(sroi, Inches(after_x), Inches(y + 0.1),
                             Inches(after_w), Inches(row_h - 0.15),
                             after_text, font_size=10, color=BLACK)

                # Divider line below each row
                add_rect(sroi, Inches(label_x), Inches(y + row_h - 0.01),
                         Inches(row_line_w), Inches(0.012), BORDER)
        # Case studies
        cases = d13.get('case_studies', [])
        if cases and not ba:
            for i, cs in enumerate(cases[:2]):
                y = 2.95 + i * 1.2
                cc = add_rect(sroi, Inches(0.55), Inches(y), Inches(8.9), Inches(1.05), LIGHT_BG)
                cc.line.color.rgb = BORDER
                cc.line.width = Pt(0.75)
                add_rect(sroi, Inches(0.55), Inches(y), Inches(0.08), Inches(1.05), ACCENT)
                add_text_box(sroi, Inches(0.8), Inches(y+0.12), Inches(3), Inches(0.35),
                             strip_emoji(cs.get('client','')), font_size=13, bold=True, color=C1, font_name=FONT_HEAD)
                add_text_box(sroi, Inches(0.8), Inches(y+0.52), Inches(8.4), Inches(0.5),
                             strip_emoji(cs.get('result','')), font_size=10, color=GRAY)
        add_footer(sroi)

    # s14: Roadmap
    # Layout-14 reference: "Roadmap" center title, 5 columns with vertical divider
    # lines. Each column shows a week header, then a cascading-down soft-tinted
    # rounded box with step title and description.
    if 's14' in sel:
        d14 = proposal.get('s14', {})
        srd = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(srd, WHITE)

        # Center title
        add_text_box(srd, Inches(0.55), Inches(0.35), Inches(8.9), Inches(0.7),
                     strip_emoji(d14.get('headline', 'Roadmap')),
                     font_size=30, bold=True, color=BLACK,
                     align=PP_ALIGN.CENTER, font_name=FONT_HEAD)
        _lead14 = strip_emoji(d14.get('lead', ''))
        if _lead14:
            add_text_box(srd, Inches(0.55), Inches(1.0), Inches(8.9), Inches(0.4),
                         _lead14, font_size=11, italic=True, color=GRAY,
                         align=PP_ALIGN.CENTER, font_name=FONT_BODY)

        milestones = d14.get('milestones', [])
        n_ms = max(min(len(milestones), 5), 1)
        col_w = 8.9 / n_ms
        base_y = 1.55
        # Vertical divider lines between columns
        for i in range(1, n_ms):
            add_rect(srd, Inches(0.55 + i * col_w - 0.01),
                     Inches(base_y + 0.3), Inches(0.02),
                     Inches(3.2), BORDER)

        # Top week headers
        for i, ms in enumerate(milestones[:5]):
            x = 0.55 + i * col_w
            add_text_box(srd, Inches(x), Inches(base_y), Inches(col_w), Inches(0.4),
                         strip_emoji(ms.get('week', '')),
                         font_size=13, bold=True, color=BLACK,
                         align=PP_ALIGN.CENTER, font_name=FONT_HEAD)

        # Cascading step cards (each row drops lower than the previous)
        card_colors = [C1_TINT, C2_TINT, C3_TINT, C4_TINT, C1_TINT]
        for i, ms in enumerate(milestones[:5]):
            col_x = 0.55 + i * col_w + 0.1
            cy = base_y + 0.6 + i * 0.55
            card_h = 0.95
            card = add_rounded_rect(srd, Inches(col_x), Inches(cy),
                                    Inches(col_w - 0.2), Inches(card_h),
                                    card_colors[i % len(card_colors)],
                                    corner_radius=Inches(0.08))
            card.line.color.rgb = C1
            card.line.width = Pt(0.75)
            add_text_box(srd, Inches(col_x + 0.1), Inches(cy + 0.08), Inches(col_w - 0.4), Inches(0.3),
                         strip_emoji(ms.get('title', '')),
                         font_size=10, bold=True, color=BLACK,
                         align=PP_ALIGN.CENTER, font_name=FONT_HEAD)
            add_text_box(srd, Inches(col_x + 0.1), Inches(cy + 0.38), Inches(col_w - 0.4), Inches(0.55),
                         strip_emoji(ms.get('desc', '')),
                         font_size=8, color=BLACK,
                         align=PP_ALIGN.CENTER)
        add_footer(srd)

    # s16: Promo Materials — matches layout-16 (3 cards, light-blue asset slot on top, multi-color status pills)
    if 's16' in sel:
        d16 = proposal.get('s16', {})
        spm = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(spm, WHITE)
        add_title_lead(spm,
                       d16.get('headline','Promo Materials'),
                       d16.get('lead',''),
                       top=0.3, left=0.55, width=9.0)
        mats = d16.get('materials', [])
        n_mats = max(len(mats[:3]), 1)
        gap = 0.25
        total_w = 10 - 1.1
        mat_w = (total_w - gap * (n_mats - 1)) / n_mats
        # Color cycle for each card's accent treatments
        pill_cycle = [C3, C2, C4]   # ensures C3 (yellow) and C4 (red) appear
        tint_cycle = [C1_TINT, C1_TINT, C1_TINT]  # asset slots all light-blue mascot slots
        top_bar_cycle = [C1, C3, C4]
        for i, mt in enumerate(mats[:3]):
            x = 0.55 + i * (mat_w + gap)
            card_top = 1.4
            card_h = 4.2
            card = add_rounded_rect(spm, Inches(x), Inches(card_top), Inches(mat_w), Inches(card_h),
                                    WHITE, corner_radius=Inches(0.12))
            card.line.color.rgb = BORDER
            card.line.width = Pt(0.75)
            # Colored top accent bar
            top_bar_c = top_bar_cycle[i % len(top_bar_cycle)]
            add_rounded_rect(spm, Inches(x), Inches(card_top), Inches(mat_w), Inches(0.14),
                             top_bar_c, corner_radius=Inches(0.06))
            # Light-blue asset slot (mascot placement zone)
            slot_m = 0.25
            slot_x = x + slot_m
            slot_y = card_top + 0.32
            slot_w = mat_w - 2 * slot_m
            slot_h = 1.85
            slot_rect = add_rounded_rect(spm, Inches(slot_x), Inches(slot_y),
                                         Inches(slot_w), Inches(slot_h),
                                         MASCOT_SLOT, corner_radius=Inches(0.08))
            slot_rect.line.color.rgb = C1
            slot_rect.line.width = Pt(1.0)
            try:
                from pptx.enum.dml import MSO_LINE_DASH_STYLE
                slot_rect.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            except Exception:
                pass
            add_text_box(spm, Inches(slot_x), Inches(slot_y + slot_h/2 - 0.15),
                         Inches(slot_w), Inches(0.3),
                         strip_emoji(mt.get('type','Asset')),
                         font_size=10, italic=True, color=C1,
                         align=PP_ALIGN.CENTER, font_name=FONT_BODY)
            # Type label
            type_y = slot_y + slot_h + 0.18
            add_text_box(spm, Inches(x+slot_m), Inches(type_y), Inches(slot_w), Inches(0.4),
                         strip_emoji(mt.get('type','')), font_size=15, bold=True,
                         color=C1, font_name=FONT_HEAD)
            # Status pill — cycles through C3/C2/C4 for "Included", gray for "Add-on"
            status = mt.get('status', '')
            if status:
                is_inc = status.strip().lower() == 'included'
                pill_color = pill_cycle[i % len(pill_cycle)] if is_inc else GRAY
                add_rounded_rect(spm, Inches(x+slot_m), Inches(type_y+0.42), Inches(1.2), Inches(0.28),
                                 pill_color, corner_radius=Inches(0.08))
                add_text_box(spm, Inches(x+slot_m), Inches(type_y+0.43), Inches(1.2), Inches(0.25),
                             ('✓ ' if is_inc else '') + strip_emoji(status),
                             font_size=9, bold=True, color=WHITE,
                             align=PP_ALIGN.CENTER, font_name=FONT_HEAD)
            add_text_box(spm, Inches(x+slot_m), Inches(type_y+0.78), Inches(slot_w), Inches(0.95),
                         strip_emoji(mt.get('desc','')), font_size=10, color=GRAY)
        add_footer(spm)

    # s17: Licensing — matches layout-17 (2x2 cards with colored left bars cycling through all brand colors)
    if 's17' in sel:
        d17 = proposal.get('s17', {})
        slic = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slic, WHITE)
        add_title_lead(slic,
                       d17.get('headline','Licensing & Ownership'),
                       d17.get('lead',''),
                       top=0.3, left=0.55, width=9.0)
        lic_cards = d17.get('cards', [])
        # Cycle through all 4 brand colors for the left vertical bars
        bar_colors = [C1, C3, C4, C2]
        num_colors = [C1, C3, C4, C2]
        # Tightened vertical rhythm so note box fits within 5.625" slide height
        card_h = 1.5
        row_gap = 0.15
        cards_top = 1.4
        for i, lc in enumerate(lic_cards[:4]):
            row, col = divmod(i, 2)
            x = 0.55 + col * 4.55
            y = cards_top + row * (card_h + row_gap)
            card = add_rounded_rect(slic, Inches(x), Inches(y), Inches(4.4), Inches(card_h),
                                    WHITE, corner_radius=Inches(0.1))
            card.line.color.rgb = BORDER
            card.line.width = Pt(0.75)
            # Colored left vertical bar (cycles all 4 brand colors)
            bar_c = bar_colors[i % len(bar_colors)]
            num_c = num_colors[i % len(num_colors)]
            add_rect(slic, Inches(x), Inches(y), Inches(0.1), Inches(card_h), bar_c)
            # Number in brand color
            add_text_box(slic, Inches(x+0.22), Inches(y+0.1), Inches(0.8), Inches(0.3),
                         f'0{i+1}', font_size=13, bold=True, color=num_c, font_name=FONT_HEAD)
            add_text_box(slic, Inches(x+0.22), Inches(y+0.38), Inches(4.0), Inches(0.4),
                         strip_emoji(lc.get('title','')), font_size=14, bold=True, color=C1, font_name=FONT_HEAD)
            add_text_box(slic, Inches(x+0.22), Inches(y+0.76), Inches(4.0), Inches(card_h - 0.8),
                         strip_emoji(lc.get('desc','')), font_size=10, color=GRAY)
        if d17.get('note'):
            # cards end at cards_top + 2*card_h + row_gap = 1.4 + 3 + 0.15 = 4.55
            note_y = cards_top + 2 * card_h + row_gap + 0.15  # 4.7
            note_h = 0.42
            note_card = add_rounded_rect(slic, Inches(0.55), Inches(note_y), Inches(8.9), Inches(note_h),
                                         C3_TINT, corner_radius=Inches(0.08))
            note_card.line.color.rgb = C3
            note_card.line.width = Pt(0.75)
            add_rect(slic, Inches(0.55), Inches(note_y), Inches(0.1), Inches(note_h), C3)
            add_text_box(slic, Inches(0.8), Inches(note_y+0.08), Inches(8.5), Inches(note_h - 0.1),
                         '★  ' + strip_emoji(d17['note']), font_size=10, italic=True, color=C1)
        add_footer(slic)

    # ══════ REORDER SLIDES ══════
    # Slides were added in order: s1,s3,s5,s6,s7,s9,s12,s15,s18 (original 9), then s2,s4,s8,s10,s11,s13,s14,s16,s17 (new 9)
    # We need to reorder them to: s1,s2,s3,s4,s5,s6,s7,s8,s9,s10,s11,s12,s13,s14,s15,s16,s17,s18
    # And filter to only selected slides

    # ══════ REORDER & FILTER SLIDES ══════
    # Original 9 slides are ALWAYS built (indices 0-8): s1, s3, s5, s6, s7, s9, s12, s15, s18
    # New slides are built conditionally (indices 9+): s2, s4, s8, s10, s11, s13, s14, s16, s17
    build_order_base = ['s1','s3','s5','s6','s7','s9','s12','s15','s18']
    build_order_new = ['s2','s4','s8','s10','s11','s13','s14','s16','s17']
    # Actual build order: base always, new only if in sel
    build_order = build_order_base + [sid for sid in build_order_new if sid in sel]
    desired_order = ['s1','s2','s3','s4','s5','s6','s7','s8','s9','s10','s11','s12','s13','s14','s15','s16','s17','s18']

    total_slides = len(prs.slides)

    if total_slides > 1:
        # Map each built slide_id to its current index in the presentation
        id_to_idx = {}
        for i, sid in enumerate(build_order):
            if i < total_slides:
                id_to_idx[sid] = i

        # Target: desired order, but only include slides that are both built AND selected
        target_order = [id_to_idx[sid] for sid in desired_order if sid in id_to_idx and sid in sel]

        # Also collect indices of slides NOT in sel (to be removed)
        remove_indices = set()
        for sid, idx in id_to_idx.items():
            if sid not in sel:
                remove_indices.add(idx)

        # Reorder using XML manipulation
        from lxml import etree
        nsmap = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        sldIdLst = prs._element.find(f'.//{nsmap}sldIdLst')
        if sldIdLst is not None:
            slides_elements = list(sldIdLst)
            if len(slides_elements) == total_slides:
                # Remove all slide refs
                for el in slides_elements:
                    sldIdLst.remove(el)
                # Re-add only selected slides in desired order
                for idx in target_order:
                    if idx < len(slides_elements):
                        sldIdLst.append(slides_elements[idx])

    prs.save(output_path)
    return output_path


# ─── PDF Generation ───
def generate_pdf(proposal, client, output_path, selected_slides=None):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm, cm
    from reportlab.lib.colors import HexColor, Color, white, black, lightgrey
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                     Table, TableStyle, PageBreak, Flowable)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import datetime, glob as globmod

    # ── Register CJK font (cross-platform: macOS, Linux, Windows) ──
    cjk_font = None
    cjk_bold = None
    font_search_paths = [
        # macOS system fonts (Japanese)
        '/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc',
        '/System/Library/Fonts/ヒラギノ角ゴシック W6.ttc',
        '/System/Library/Fonts/Hiragino Sans GB.ttc',
        '/Library/Fonts/Arial Unicode.ttf',
        # macOS - Noto fonts if installed via Homebrew etc.
        '/Library/Fonts/NotoSansCJK-Regular.ttc',
        '/Library/Fonts/NotoSansCJKjp-Regular.otf',
        # macOS PingFang (Chinese)
        '/System/Library/Fonts/PingFang.ttc',
        '/System/Library/Fonts/STHeiti Light.ttc',
        '/System/Library/Fonts/STHeiti Medium.ttc',
        # Linux
        '/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc',
        '/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc',
        '/usr/share/fonts/truetype/droid/DroidSansFallback.ttf',
        '/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc',
        # Windows
        'C:/Windows/Fonts/msgothic.ttc',
        'C:/Windows/Fonts/msmincho.ttc',
        'C:/Windows/Fonts/meiryo.ttc',
        'C:/Windows/Fonts/YuGothR.ttc',
    ]

    # Also search common macOS font dirs with glob
    for pattern in ['/System/Library/Fonts/*.ttc', '/System/Library/Fonts/Supplemental/*.ttf',
                    '/System/Library/Fonts/Supplemental/*.ttc', '/Library/Fonts/*.ttf', '/Library/Fonts/*.ttc']:
        try:
            font_search_paths.extend(globmod.glob(pattern))
        except:
            pass

    for font_path in font_search_paths:
        if not os.path.exists(font_path):
            continue
        try:
            pdfmetrics.registerFont(TTFont('CJK', font_path, subfontIndex=0))
            cjk_font = 'CJK'
            # Try to register bold variant
            try:
                pdfmetrics.registerFont(TTFont('CJKBold', font_path, subfontIndex=0))
                cjk_bold = 'CJKBold'
            except:
                cjk_bold = cjk_font
            break
        except:
            try:
                pdfmetrics.registerFont(TTFont('CJK', font_path))
                cjk_font = 'CJK'
                cjk_bold = 'CJK'
                break
            except:
                continue

    if not cjk_font:
        # Fallback: use reportlab's built-in CIDFont for CJK
        try:
            from reportlab.pdfbase.cidfonts import UnicodeCIDFont
            pdfmetrics.registerFont(UnicodeCIDFont('HeiseiKakuGo-W5'))
            cjk_font = 'HeiseiKakuGo-W5'
            cjk_bold = 'HeiseiKakuGo-W5'
        except:
            try:
                from reportlab.pdfbase.cidfonts import UnicodeCIDFont
                pdfmetrics.registerFont(UnicodeCIDFont('STSong-Light'))
                cjk_font = 'STSong-Light'
                cjk_bold = 'STSong-Light'
            except:
                cjk_font = 'Helvetica'
                cjk_bold = 'Helvetica-Bold'

    c1 = client.get('color1', '#0A0A0A')
    c2 = client.get('color2', '#FFFFFF')
    c1_color = HexColor(c1)
    c2_color = HexColor(c2)

    def strip_emoji(s):
        if not s: return ''
        return re.sub(r'[\U0001F300-\U0001FAFF\U00002600-\U000027BF\U0000FE00-\U0000FE0F\U0001F900-\U0001F9FF\U0000200D\U000020E3\U000E0020-\U000E007F]', '', str(s)).strip()

    def lum(c):
        return 0.299*c.red + 0.587*c.green + 0.114*c.blue

    text_on_brand = white if lum(c1_color) < 0.5 else black
    brand_hex = c1
    page_w, page_h = A4

    # ── Custom Flowable: colored rectangle block ──
    class ColorBlock(Flowable):
        """A colored rectangle with text content inside."""
        def __init__(self, width, height, bg_color, content_paragraphs=None):
            Flowable.__init__(self)
            self.width = width
            self.height = height
            self.bg_color = bg_color
            self.content_paragraphs = content_paragraphs or []

        def draw(self):
            self.canv.setFillColor(self.bg_color)
            self.canv.roundRect(0, 0, self.width, self.height, 6, fill=1, stroke=0)
            y = self.height - 15*mm
            for p, style in self.content_paragraphs:
                pw, ph = p.wrap(self.width - 20*mm, self.height)
                p.drawOn(self.canv, 10*mm, y - ph)
                y -= (ph + 3*mm)

    class HLine(Flowable):
        """A horizontal line with brand color."""
        def __init__(self, width, color=None, thickness=0.5):
            Flowable.__init__(self)
            self.width = width
            self.line_color = color or c1_color
            self.thickness = thickness

        def draw(self):
            self.canv.setStrokeColor(self.line_color)
            self.canv.setLineWidth(self.thickness)
            self.canv.line(0, 0, self.width, 0)

        def wrap(self, aW, aH):
            return (self.width, self.thickness + 2*mm)

    # ── Styles ──
    styles = getSampleStyleSheet()
    body_font = cjk_font
    bold_font = cjk_bold or cjk_font
    usable_w = page_w - 40*mm

    styles.add(ParagraphStyle('CoverTitle', fontName=bold_font, fontSize=28, leading=36,
                               textColor=text_on_brand, spaceAfter=6))
    styles.add(ParagraphStyle('CoverSub', fontName=body_font, fontSize=13, leading=18,
                               textColor=text_on_brand, spaceAfter=4))
    styles.add(ParagraphStyle('CoverQuote', fontName=body_font, fontSize=11, leading=16,
                               textColor=text_on_brand, spaceAfter=4, leftIndent=10*mm))
    styles.add(ParagraphStyle('SectionHead', fontName=bold_font, fontSize=16, leading=22,
                               textColor=c1_color, spaceBefore=10, spaceAfter=6))
    styles.add(ParagraphStyle('SubHead', fontName=bold_font, fontSize=12, leading=16,
                               textColor=HexColor('#1e1e1e'), spaceBefore=8, spaceAfter=4))
    styles.add(ParagraphStyle('BodyText2', fontName=body_font, fontSize=10, leading=15,
                               textColor=HexColor('#444444'), spaceAfter=6))
    styles.add(ParagraphStyle('BoldText', fontName=bold_font, fontSize=10, leading=15,
                               textColor=HexColor('#1e1e1e'), spaceAfter=3))
    styles.add(ParagraphStyle('SmallGray', fontName=body_font, fontSize=9, leading=12,
                               textColor=HexColor('#888888'), spaceAfter=4))
    styles.add(ParagraphStyle('ChatUser', fontName=bold_font, fontSize=10, leading=15,
                               textColor=HexColor('#333333'), spaceAfter=2, leftIndent=5*mm))
    styles.add(ParagraphStyle('ChatBot', fontName=body_font, fontSize=10, leading=15,
                               textColor=c1_color, spaceAfter=6, leftIndent=5*mm))
    styles.add(ParagraphStyle('StatNum', fontName=bold_font, fontSize=14, leading=18,
                               textColor=c1_color, spaceAfter=1))
    styles.add(ParagraphStyle('StatLabel', fontName=body_font, fontSize=9, leading=12,
                               textColor=HexColor('#666666'), spaceAfter=6))
    styles.add(ParagraphStyle('TierName', fontName=bold_font, fontSize=13, leading=18,
                               textColor=HexColor('#1e1e1e'), spaceAfter=2))
    styles.add(ParagraphStyle('TierPrice', fontName=bold_font, fontSize=16, leading=22,
                               textColor=c1_color, spaceAfter=4))
    styles.add(ParagraphStyle('FooterText', fontName=body_font, fontSize=8, leading=10,
                               textColor=HexColor('#AAAAAA'), alignment=TA_CENTER))

    # ── Page background with brand header ──
    def on_first_page(canvas, doc):
        """Cover page: full brand color background."""
        canvas.saveState()
        canvas.setFillColor(c1_color)
        canvas.rect(0, 0, page_w, page_h, fill=1, stroke=0)
        # Bottom accent line
        canvas.setFillColor(c2_color)
        canvas.rect(0, 0, page_w, 3*mm, fill=1, stroke=0)
        canvas.restoreState()

    def on_later_pages(canvas, doc):
        """Content pages: white bg with brand color header bar + footer."""
        canvas.saveState()
        # Top brand bar
        canvas.setFillColor(c1_color)
        canvas.rect(0, page_h - 8*mm, page_w, 8*mm, fill=1, stroke=0)
        # Footer line
        canvas.setStrokeColor(HexColor('#E0E0E0'))
        canvas.setLineWidth(0.5)
        canvas.line(20*mm, 15*mm, page_w - 20*mm, 15*mm)
        # Footer text
        canvas.setFont(body_font, 7)
        canvas.setFillColor(HexColor('#BBBBBB'))
        canvas.drawCentredString(page_w/2, 10*mm, 'Powered by notso.ai')
        canvas.restoreState()

    doc = SimpleDocTemplate(output_path, pagesize=A4,
                            leftMargin=20*mm, rightMargin=20*mm,
                            topMargin=25*mm, bottomMargin=22*mm)

    # Selected slides filter
    sel = set(selected_slides) if selected_slides else {'s1','s2','s3','s4','s5','s6','s7','s8','s9','s10','s11','s12','s13','s14','s15','s16','s17','s18'}

    story = []
    section_num = 0

    def next_section():
        nonlocal section_num
        section_num += 1
        return section_num

    # ══════ s1: COVER PAGE ══════
    if 's1' in sel:
        d1 = proposal.get('s1',{})
        story.append(Spacer(1, 50*mm))
        story.append(Paragraph(strip_emoji(client.get('name','Proposal')), styles['CoverTitle']))
        story.append(Paragraph(strip_emoji(d1.get('tagline','')), styles['CoverSub']))
        story.append(Spacer(1, 15*mm))
        story.append(Paragraph('"' + strip_emoji(d1.get('greeting','')) + '"', styles['CoverQuote']))
        story.append(Spacer(1, 4*mm))
        story.append(Paragraph('— ' + strip_emoji(d1.get('mascot_name','')) + ', powered by notso.ai', styles['CoverQuote']))
        story.append(Spacer(1, 30*mm))
        story.append(Paragraph(strip_emoji(client.get('industry','')) + '  |  ' + datetime.datetime.now().strftime('%d %B %Y'), styles['CoverSub']))
        story.append(PageBreak())

    # ══════ s2: TABLE OF CONTENTS ══════
    if 's2' in sel:
        n = next_section()
        toc_labels = {'s1':'Cover','s3':'Pain Points','s4':'Market Opportunity','s5':'Core Features',
                      's6':'Mascot Selection','s7':'Mascot Design','s8':'Personality & Empathy',
                      's9':'Chat Mock-up','s10':'Chatflow Design','s11':'Knowledge Base',
                      's12':'Data & Insights','s13':'ROI Evidence','s14':'Roadmap',
                      's15':'Pricing','s16':'Promo Materials','s17':'Licensing','s18':'Thank You'}
        story.append(Paragraph(f'{n}. Agenda', styles['SectionHead']))
        story.append(HLine(usable_w))
        story.append(Spacer(1, 3*mm))
        toc_items = [s for s in ['s1','s3','s4','s5','s6','s7','s8','s9','s10','s11','s12','s13','s14','s15','s16','s17','s18'] if s in sel]
        for i, sid in enumerate(toc_items):
            story.append(Paragraph(f'{i+1:02d}  {toc_labels.get(sid, sid)}', styles['BodyText2']))
        story.append(PageBreak())

    # ══════ s3: PAIN POINTS ══════
    if 's3' in sel:
        n = next_section()
        d3 = proposal.get('s3',{})
        story.append(Paragraph(f'{n}. ' + strip_emoji(d3.get('headline','Pain Points')), styles['SectionHead']))
        story.append(HLine(usable_w))
        story.append(Paragraph(strip_emoji(d3.get('intro','')), styles['BodyText2']))
        story.append(Spacer(1, 3*mm))
        for p in d3.get('points',[]):
            story.append(Paragraph(strip_emoji(p.get('title','')), styles['SubHead']))
            story.append(Paragraph(strip_emoji(p.get('desc','')), styles['BodyText2']))
        story.append(PageBreak())

    # ══════ s4: MARKET OPPORTUNITY ══════
    if 's4' in sel:
        n = next_section()
        d4 = proposal.get('s4',{})
        story.append(Paragraph(f'{n}. ' + strip_emoji(d4.get('headline','Market Opportunity')), styles['SectionHead']))
        story.append(HLine(usable_w))
        story.append(Spacer(1, 3*mm))
        for label, key in [('Industry Size','industry_size'),('Growth Rate','growth_rate'),('Projected Size','projected_size')]:
            data_val = d4.get(key, {})
            if data_val:
                story.append(Paragraph(f'{label}: {strip_emoji(str(data_val.get("value","")))}', styles['SubHead']))
                story.append(Paragraph('Source: ' + strip_emoji(str(data_val.get('source',''))), styles['SmallGray']))
        comps = d4.get('competitors', [])
        if comps:
            story.append(Spacer(1, 4*mm))
            story.append(Paragraph('Competitive Gaps', styles['SubHead']))
            for c in comps[:4]:
                story.append(Paragraph(f'{strip_emoji(c.get("name",""))}: {strip_emoji(c.get("gap",""))}', styles['BodyText2']))
        story.append(PageBreak())

    # ══════ s5: FEATURES ══════
    if 's5' in sel:
        n = next_section()
        d5 = proposal.get('s5',{})
        story.append(Paragraph(f'{n}. ' + strip_emoji(d5.get('headline','Core Features')), styles['SectionHead']))
        story.append(HLine(usable_w))
        story.append(Paragraph(strip_emoji(d5.get('intro','')), styles['BodyText2']))
        story.append(Spacer(1, 3*mm))
        for f in d5.get('features',[]):
            story.append(Paragraph(strip_emoji(f.get('title','')), styles['SubHead']))
            story.append(Paragraph(strip_emoji(f.get('desc','')), styles['BodyText2']))
        story.append(PageBreak())

    # ══════ s6: MASCOT OPTIONS ══════
    if 's6' in sel:
        n = next_section()
        d6 = proposal.get('s6',{})
        story.append(Paragraph(f'{n}. ' + strip_emoji(d6.get('headline','Mascot Selection')), styles['SectionHead']))
        story.append(HLine(usable_w))
        for key, label in [('option_a','A'), ('option_b','B'), ('option_c','C')]:
            opt = d6.get(key)
            if opt:
                story.append(Paragraph(f'Option {label}: {strip_emoji(opt.get("name",""))}', styles['SubHead']))
                story.append(Paragraph(strip_emoji(opt.get('desc','')), styles['BodyText2']))
        story.append(PageBreak())

    # ══════ s7: PERSONALITY ══════
    if 's7' in sel:
        n = next_section()
        d7 = proposal.get('s7',{})
        story.append(Paragraph(f'{n}. Personality — ' + strip_emoji(d7.get('name','')), styles['SectionHead']))
        story.append(HLine(usable_w))
        story.append(Paragraph(strip_emoji(d7.get('personality','')), styles['BodyText2']))
        story.append(Paragraph('Tone: ' + strip_emoji(d7.get('tone_desc','')), styles['BodyText2']))
        phrases = d7.get('phrases',[])
        if phrases:
            story.append(Paragraph(' | '.join(['"'+strip_emoji(p)+'"' for p in phrases]), styles['SmallGray']))
        story.append(PageBreak())

    # ══════ s8: PERSONALITY & EMPATHY (Expression Grid) ══════
    if 's8' in sel:
        n = next_section()
        d8 = proposal.get('s8',{})
        story.append(Paragraph(f'{n}. Personality & Empathy — ' + strip_emoji(d8.get('name','')), styles['SectionHead']))
        story.append(HLine(usable_w))
        story.append(Paragraph(strip_emoji(d8.get('personality_summary','')), styles['BodyText2']))
        story.append(Spacer(1, 3*mm))
        expressions = d8.get('expressions', [])
        for expr in expressions[:9]:
            story.append(Paragraph(f'{strip_emoji(expr.get("emotion",""))}: {strip_emoji(expr.get("context",""))}', styles['BodyText2']))
        story.append(PageBreak())

    # ══════ s9: CHAT DEMO ══════
    if 's9' in sel:
        n = next_section()
        d9 = proposal.get('s9',{})
        story.append(Paragraph(f'{n}. Chat Demo — ' + strip_emoji(d9.get('mascot_name','')), styles['SectionHead']))
        story.append(HLine(usable_w))
        story.append(Spacer(1, 3*mm))
        for m in d9.get('chat',[]):
            is_user = m.get('r') == 'user'
            label = 'User' if is_user else strip_emoji(d9.get('mascot_name','Bot'))
            style = styles['ChatUser'] if is_user else styles['ChatBot']
            story.append(Paragraph(f'{label}: {strip_emoji(m.get("m",""))}', style))
            story.append(Spacer(1, 2*mm))
        story.append(PageBreak())

    # ══════ s10: CHATFLOW DESIGN ══════
    if 's10' in sel:
        n = next_section()
        d10 = proposal.get('s10',{})
        story.append(Paragraph(f'{n}. ' + strip_emoji(d10.get('headline','How the conversation flows')), styles['SectionHead']))
        story.append(HLine(usable_w))
        story.append(Spacer(1, 3*mm))
        for i, step in enumerate(d10.get('steps',[])):
            story.append(Paragraph(f'Step {i+1}: {strip_emoji(step.get("title",""))}', styles['SubHead']))
            story.append(Paragraph(strip_emoji(step.get('desc','')), styles['BodyText2']))
        story.append(PageBreak())

    # ══════ s11: KNOWLEDGE BASE ══════
    if 's11' in sel:
        n = next_section()
        d11 = proposal.get('s11',{})
        story.append(Paragraph(f'{n}. ' + strip_emoji(d11.get('headline','Knowledge Base')), styles['SectionHead']))
        story.append(HLine(usable_w))
        story.append(Spacer(1, 3*mm))
        for cat in d11.get('categories',[]):
            story.append(Paragraph(strip_emoji(cat.get('title','')), styles['SubHead']))
            for it in cat.get('items',[]):
                story.append(Paragraph(f'  {strip_emoji(it)}', styles['BodyText2']))
            docs = cat.get('docs',[])
            if docs:
                story.append(Paragraph('Documents: ' + ', '.join([strip_emoji(d) for d in docs]), styles['SmallGray']))
        story.append(PageBreak())

    # ══════ s12: DATA & INSIGHTS ══════
    if 's12' in sel:
        n = next_section()
        d12 = proposal.get('s12',{})
        story.append(Paragraph(f'{n}. ' + strip_emoji(d12.get('headline','Data & Insights')), styles['SectionHead']))
        story.append(HLine(usable_w))
        # Support both old format (stats with n/l) and new format (metrics with icon/label/desc)
        metrics = d12.get('metrics', d12.get('stats', []))
        if metrics:
            stat_data = []
            for st in metrics:
                num = strip_emoji(st.get('n', st.get('icon', '')))
                lbl = strip_emoji(st.get('l', st.get('label', '')))
                desc = strip_emoji(st.get('desc', ''))
                stat_data.append([
                    Paragraph(num, styles['StatNum']),
                    Paragraph(lbl + (' — ' + desc if desc else ''), styles['StatLabel'])
                ])
            if stat_data:
                col_w = usable_w / 2
                st_table = Table(stat_data, colWidths=[col_w*0.3, col_w*1.7])
                st_table.setStyle(TableStyle([
                    ('VALIGN', (0,0), (-1,-1), 'TOP'),
                    ('TOPPADDING', (0,0), (-1,-1), 2),
                    ('BOTTOMPADDING', (0,0), (-1,-1), 2),
                ]))
                story.append(st_table)
        badges = d12.get('badges', [])
        if badges:
            story.append(Spacer(1, 3*mm))
            badge_text = ' | '.join([strip_emoji(b) for b in badges])
            story.append(Paragraph(badge_text, styles['SmallGray']))
        if d12.get('case_client'):
            story.append(Spacer(1, 4*mm))
            story.append(Paragraph(f'Case Study: {strip_emoji(d12["case_client"])}', styles['SubHead']))
            story.append(Paragraph(strip_emoji(d12.get('case_result','')), styles['BodyText2']))
        story.append(PageBreak())

    # ══════ s13: ROI EVIDENCE ══════
    if 's13' in sel:
        n = next_section()
        d13 = proposal.get('s13',{})
        story.append(Paragraph(f'{n}. ' + strip_emoji(d13.get('headline','Proven ROI')), styles['SectionHead']))
        story.append(HLine(usable_w))
        roi_stats = d13.get('stats', [])
        if roi_stats:
            stat_data = []
            for st in roi_stats:
                stat_data.append([
                    Paragraph(strip_emoji(st.get('n','')), styles['StatNum']),
                    Paragraph(strip_emoji(st.get('l','')) + (' — ' + strip_emoji(st.get('detail','')) if st.get('detail') else ''), styles['StatLabel'])
                ])
            col_w = usable_w / 2
            st_table = Table(stat_data, colWidths=[col_w*0.3, col_w*1.7])
            st_table.setStyle(TableStyle([('VALIGN',(0,0),(-1,-1),'TOP'),('TOPPADDING',(0,0),(-1,-1),2),('BOTTOMPADDING',(0,0),(-1,-1),2)]))
            story.append(st_table)
        ba = d13.get('before_after', {})
        if ba:
            story.append(Spacer(1, 4*mm))
            story.append(Paragraph('Before', styles['SubHead']))
            for b in ba.get('before',[]):
                story.append(Paragraph(f'  {strip_emoji(b)}', styles['BodyText2']))
            story.append(Paragraph('After notso.ai', styles['SubHead']))
            for a in ba.get('after',[]):
                story.append(Paragraph(f'  {strip_emoji(a)}', styles['BodyText2']))
        cases = d13.get('case_studies', [])
        if cases:
            for cs in cases[:3]:
                story.append(Paragraph(strip_emoji(cs.get('client','')), styles['SubHead']))
                story.append(Paragraph(strip_emoji(cs.get('result','')), styles['BodyText2']))
        story.append(PageBreak())

    # ══════ s14: ROADMAP ══════
    if 's14' in sel:
        n = next_section()
        d14 = proposal.get('s14',{})
        story.append(Paragraph(f'{n}. ' + strip_emoji(d14.get('headline','Your roadmap to launch')), styles['SectionHead']))
        story.append(HLine(usable_w))
        story.append(Spacer(1, 3*mm))
        for ms in d14.get('milestones',[]):
            story.append(Paragraph(f'{strip_emoji(ms.get("week",""))}: {strip_emoji(ms.get("title",""))}', styles['SubHead']))
            story.append(Paragraph(strip_emoji(ms.get('desc','')), styles['BodyText2']))
        story.append(PageBreak())

    # ══════ s15: PRICING ══════
    if 's15' in sel:
        n = next_section()
        d15 = proposal.get('s15',{})
        story.append(Paragraph(f'{n}. Pricing', styles['SectionHead']))
        story.append(HLine(usable_w))
        if d15.get('reasoning'):
            story.append(Paragraph(strip_emoji(d15['reasoning']), styles['BodyText2']))
        story.append(Spacer(1, 4*mm))
        tiers = d15.get('tiers', d15.get('plans', []))
        if tiers:
            tier_rows = []
            for t in tiers:
                rec_tier = d15.get('rec_tier', '')
                is_rec = rec_tier and t.get('name','').lower().find(rec_tier.lower()) >= 0
                rec_tag = ' [RECOMMENDED]' if is_rec else ''
                features_text = ''
                if t.get('features'):
                    features_text = ', '.join([strip_emoji(f) for f in t['features'][:5]])
                elif t.get('highlight'):
                    features_text = strip_emoji(t['highlight'])
                tier_rows.append([
                    Paragraph(f'{strip_emoji(t.get("name",""))}{rec_tag}', styles['TierName']),
                    Paragraph(strip_emoji(t.get('price','')), styles['TierPrice']),
                    Paragraph(features_text, styles['BodyText2']),
                ])
            tier_table = Table(tier_rows, colWidths=[usable_w*0.25, usable_w*0.25, usable_w*0.5])
            tier_style = [
                ('VALIGN', (0,0), (-1,-1), 'TOP'),
                ('TOPPADDING', (0,0), (-1,-1), 6),
                ('BOTTOMPADDING', (0,0), (-1,-1), 6),
                ('LINEBELOW', (0,0), (-1,-2), 0.5, HexColor('#E0E0E0')),
            ]
            tier_table.setStyle(TableStyle(tier_style))
            story.append(tier_table)
        addons = d15.get('addons', [])
        if addons:
            story.append(Spacer(1, 4*mm))
            story.append(Paragraph('Add-ons', styles['SubHead']))
            for addon in addons:
                story.append(Paragraph(f'{strip_emoji(addon.get("name",""))}: {strip_emoji(addon.get("price",""))}', styles['BodyText2']))
        if d15.get('next_step'):
            story.append(Spacer(1, 4*mm))
            story.append(Paragraph(strip_emoji(d15['next_step']), styles['BoldText']))
        story.append(PageBreak())

    # ══════ s16: PROMO MATERIALS ══════
    if 's16' in sel:
        n = next_section()
        d16 = proposal.get('s16',{})
        story.append(Paragraph(f'{n}. ' + strip_emoji(d16.get('headline','Promo Materials')), styles['SectionHead']))
        story.append(HLine(usable_w))
        story.append(Spacer(1, 3*mm))
        for mt in d16.get('materials',[]):
            status = mt.get('status','')
            prefix = ' [Included]' if status == 'Included' else ''
            story.append(Paragraph(f'{strip_emoji(mt.get("type",""))}{prefix}', styles['SubHead']))
            story.append(Paragraph(strip_emoji(mt.get('desc','')), styles['BodyText2']))
        story.append(PageBreak())

    # ══════ s17: LICENSING ══════
    if 's17' in sel:
        n = next_section()
        d17 = proposal.get('s17',{})
        story.append(Paragraph(f'{n}. ' + strip_emoji(d17.get('headline','Licensing & Ownership')), styles['SectionHead']))
        story.append(HLine(usable_w))
        story.append(Spacer(1, 3*mm))
        for lc in d17.get('cards',[]):
            story.append(Paragraph(strip_emoji(lc.get('title','')), styles['SubHead']))
            story.append(Paragraph(strip_emoji(lc.get('desc','')), styles['BodyText2']))
        if d17.get('note'):
            story.append(Spacer(1, 4*mm))
            story.append(Paragraph(strip_emoji(d17['note']), styles['SmallGray']))
        story.append(PageBreak())

    # ══════ s18: THANK YOU ══════
    if 's18' in sel:
        n = next_section()
        d18 = proposal.get('s18',{})
        story.append(Paragraph(f'{n}. Next Steps', styles['SectionHead']))
        story.append(HLine(usable_w))
        story.append(Paragraph(strip_emoji(d18.get('closing','')), styles['BodyText2']))
        story.append(Spacer(1, 3*mm))
        for i, s in enumerate(d18.get('next_steps',[])):
            story.append(Paragraph(f'{i+1}. {strip_emoji(s)}', styles['BodyText2']))
        story.append(Spacer(1, 10*mm))
        contact = ' | '.join(filter(None, [
            d18.get('email','hello@notso.ai'),
            d18.get('phone',''),
            d18.get('website','www.notso.ai')
        ]))
        story.append(Paragraph('notso.ai', styles['SubHead']))
        story.append(Paragraph(contact, styles['SmallGray']))

    doc.build(story, onFirstPage=on_first_page, onLaterPages=on_later_pages)
    return output_path


# ─── Main ───
if __name__ == '__main__':
    data = json.loads(sys.stdin.read())
    fmt = data.get('format', 'pptx')
    proposal = data.get('proposal', {})
    client_data = data.get('client', {})
    output_dir = data.get('output_dir', '/tmp')
    selected_slides = data.get('selected_slides', None)

    name = re.sub(r'[^\w\-]', '', (client_data.get('name') or 'draft').replace(' ', '-'))

    if fmt == 'pptx':
        path = os.path.join(output_dir, f'notso-proposal-{name}.pptx')
        generate_pptx(proposal, client_data, path, selected_slides=selected_slides)
    elif fmt == 'pdf':
        path = os.path.join(output_dir, f'notso-proposal-{name}.pdf')
        generate_pdf(proposal, client_data, path, selected_slides=selected_slides)
    else:
        print(json.dumps({'error': f'Unknown format: {fmt}'}))
        sys.exit(1)

    print(json.dumps({'success': True, 'path': path, 'filename': os.path.basename(path)}))
